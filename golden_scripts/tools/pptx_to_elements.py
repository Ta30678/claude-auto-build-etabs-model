"""
PPT-to-Elements Script
======================
Converts PowerPoint (.pptx) structural plan drawings into structural element
JSON (columns, beams, walls, small beams) with precise coordinates.

PPT format: FREEFORM shapes on top of PNG base images, with color-coded legend.

Usage:
    # Phase 1: major beams + columns + walls
    python -m golden_scripts.tools.pptx_to_elements \
        --input plan.pptx --output elements.json \
        --page-floors "1=B3F, 3=1F~2F, 4=3F~14F" \
        --phase phase1

    # Phase 2: small beams only
    python -m golden_scripts.tools.pptx_to_elements \
        --input plan.pptx --output sb_elements.json \
        --page-floors "3=1F~2F, 4=3F~14F" \
        --phase phase2

    # Scan slides for floor labels (auto-detect)
    python -m golden_scripts.tools.pptx_to_elements --input plan.pptx --scan-floors

    # Auto-detect floors and process (no --page-floors needed)
    python -m golden_scripts.tools.pptx_to_elements \
        --input plan.pptx --output elements.json \
        --auto-floors --phase phase1

    # List slides and their shape counts
    python -m golden_scripts.tools.pptx_to_elements --input plan.pptx --list-slides

    # Extract + crop PNGs
    python -m golden_scripts.tools.pptx_to_elements \
        --input plan.pptx --output elements.json \
        --page-floors "1=B3F" --phase phase1 \
        --crop --crop-dir "./結構配置圖/"

    # Preview without writing
    python -m golden_scripts.tools.pptx_to_elements ... --dry-run
"""

import argparse
import json
import math
import re
import sys
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Optional

from lxml import etree
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR_INDEX
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu


# ─── Shared Constants & Helpers ──────────────────────────────────────────────
# (Previously imported from annot_to_elements.py, now self-contained)

ROUND_5CM = 5                   # Round column dimensions to nearest 5cm
COORD_PRECISION = 2             # decimal places (0.01m = 1cm)
DIRECTION_RATIO = 0.1           # dx/dy or dy/dx threshold for axis detection

# Regex for section names like SB30x60, FB100x230, FWB60x230, B55x80
_SECTION_RE = re.compile(
    r"(FWB|FSB|FB|WB|SB|B)(\d+)[xX](\d+)", re.IGNORECASE
)

# Regex for wall thickness like "90CM 連續壁" or "25cm壁"
_WALL_CM_RE = re.compile(r"(\d+)\s*[cC][mM]\s*(?:連續壁|壁)")

# Floor label patterns for --scan-floors / --auto-floors
# Word boundaries prevent false positives (e.g. C490F should NOT match)
_FLOOR_WB_BEFORE = r"(?<![A-EG-QS-Z\d])"  # Not preceded by letters (except F,R) or digits
_FLOOR_WB_AFTER = r"(?![A-Za-z0-9])"        # Not followed by ASCII alphanumeric
_FLOOR_SINGLE_RE = rf"{_FLOOR_WB_BEFORE}[BR]?\d+F{_FLOOR_WB_AFTER}"
_FLOOR_RANGE_RE = rf"{_FLOOR_WB_BEFORE}[BR]?\d+F\s*[~\-\uff5e]\s*[BR]?\d+F{_FLOOR_WB_AFTER}"
FLOOR_LABEL_RE = re.compile(
    rf"({_FLOOR_RANGE_RE}|{_FLOOR_SINGLE_RE})", re.IGNORECASE
)


@dataclass
class LegendEntry:
    """A legend item mapping a color to a structural element type."""
    element_type: str           # "beam", "small_beam", "column", "wall"
    section: Optional[str]      # e.g. "B55X80", "SB30X60", or None (generic)
    color_name: str
    color_rgb: list
    specificity: int            # 0=generic, 1=specific
    is_diaphragm: bool = False
    prefix: str = ""            # B, SB, WB, FB, FWB, FSB, C, W
    label: str = ""             # original label text


def expand_floor_range(range_str: str) -> list[str]:
    """Expand floor range like '3F~14F' into ['3F', '4F', ..., '14F'].

    Supported patterns:
      B3F~B1F   → [B3F, B2F, B1F]      (basement descending)
      B3F~1F    → [B3F, B2F, B1F, 1F]   (basement to ground)
      1F~14F    → [1F, 2F, ..., 14F]     (normal ascending)
      R1F~R3F   → [R1F, R2F, R3F]        (rooftop ascending)
      RF        → [RF]                    (single floor)
    """
    range_str = range_str.strip()
    range_str = range_str.replace("\uff5e", "~").replace("-", "~")  # normalize separators
    if "~" not in range_str:
        return [range_str]

    start, end = [s.strip() for s in range_str.split("~", 1)]

    # Basement: B3F~B1F
    ms = re.match(r"^B(\d+)F$", start)
    me = re.match(r"^B(\d+)F$", end)
    if ms and me:
        s, e = int(ms.group(1)), int(me.group(1))
        step = -1 if s > e else 1
        return [f"B{i}F" for i in range(s, e + step, step)]

    # Basement to ground: B3F~1F  or  B3F~2F
    if ms and re.match(r"^(\d+)F$", end):
        s = int(ms.group(1))
        e = int(re.match(r"^(\d+)F$", end).group(1))
        result = [f"B{i}F" for i in range(s, 0, -1)]
        result.extend(f"{i}F" for i in range(1, e + 1))
        return result

    # Normal floors: 3F~14F
    ms = re.match(r"^(\d+)F$", start)
    me = re.match(r"^(\d+)F$", end)
    if ms and me:
        s, e = int(ms.group(1)), int(me.group(1))
        return [f"{i}F" for i in range(s, e + 1)]

    # Rooftop: R1F~R3F
    ms = re.match(r"^R(\d+)F$", start)
    me = re.match(r"^R(\d+)F$", end)
    if ms and me:
        s, e = int(ms.group(1)), int(me.group(1))
        return [f"R{i}F" for i in range(s, e + 1)]

    # Normal floor to RF: 12F~RF (we cannot enumerate without max floor)
    # Normal floor to R1F:  similar
    ms_n = re.match(r"^(\d+)F$", start)
    me_r = re.match(r"^R(\d*)F$", end)
    if ms_n and me_r:
        s = int(ms_n.group(1))
        r_num = me_r.group(1)
        result = [f"{i}F" for i in range(s, s + 1)]  # start only
        # We'll append RF / R1F etc. — caller must provide full floor list if needed
        result.append(end)
        print(f"  WARNING: Cannot fully enumerate '{range_str}'; returning [{start}, {end}]")
        return result

    # Fallback
    print(f"  WARNING: Cannot expand floor range '{range_str}'; returning as-is")
    return [start, end]


def parse_page_floors(page_floors_str: str) -> dict[int, list[str]]:
    """Parse --page-floors argument.

    Format: "1=B3F, 3=1F~2F, 4=3F~14F, 5=R1F~R3F"
    Returns: {1: ["B3F"], 3: ["1F","2F"], 4: ["3F",..."14F"], ...}
    """
    result = {}
    for item in page_floors_str.split(","):
        item = item.strip()
        if "=" not in item:
            continue
        page_str, floor_range = item.split("=", 1)
        page_num = int(page_str.strip())
        floors = expand_floor_range(floor_range.strip())
        result[page_num] = floors
    return result


class _AbsShape:
    """Proxy that overrides left/top with absolute slide coordinates."""
    __slots__ = ('_shape', 'left', 'top')

    def __init__(self, shape, abs_left, abs_top):
        self._shape = shape
        self.left = abs_left
        self.top = abs_top

    def __getattr__(self, name):
        return getattr(self._shape, name)


_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _get_group_xfrm(group_shape):
    """Extract group transform: (off_x, off_y, ext_cx, ext_cy, ch_off_x, ch_off_y, ch_ext_cx, ch_ext_cy)."""
    sp_tree = group_shape._element
    # grpSpPr/xfrm is the canonical location for group transforms
    grp_sp_pr = sp_tree.find(f'{{{_NS_A}}}../grpSpPr') if False else None
    xfrm_el = None
    # Search in the group shape's XML element for xfrm under grpSpPr (PowerPoint namespace)
    _NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    _NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    for ns_prefix in [
        f'{{{_NS_A}}}',
        '{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}',
    ]:
        xfrm_el = sp_tree.find(f'.//grpSpPr/{{{_NS_A}}}xfrm', sp_tree.nsmap if hasattr(sp_tree, 'nsmap') else {})
        if xfrm_el is not None:
            break
    if xfrm_el is None:
        # Try direct child search with lxml
        for child in sp_tree:
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag == 'grpSpPr':
                for sub in child:
                    sub_tag = sub.tag.split('}')[-1] if '}' in sub.tag else sub.tag
                    if sub_tag == 'xfrm':
                        xfrm_el = sub
                        break
                break
    if xfrm_el is None:
        # Fallback: use python-pptx properties, assume identity child transform
        return (group_shape.left, group_shape.top,
                group_shape.width, group_shape.height,
                0, 0, group_shape.width, group_shape.height)

    ns = _NS_A
    off = xfrm_el.find(f'{{{ns}}}off')
    ext = xfrm_el.find(f'{{{ns}}}ext')
    ch_off = xfrm_el.find(f'{{{ns}}}chOff')
    ch_ext = xfrm_el.find(f'{{{ns}}}chExt')

    ox  = int(off.get('x', 0))    if off is not None else 0
    oy  = int(off.get('y', 0))    if off is not None else 0
    ecx = int(ext.get('cx', 1))   if ext is not None else 1
    ecy = int(ext.get('cy', 1))   if ext is not None else 1
    cox = int(ch_off.get('x', 0)) if ch_off is not None else 0
    coy = int(ch_off.get('y', 0)) if ch_off is not None else 0
    cecx = int(ch_ext.get('cx', ecx)) if ch_ext is not None else ecx
    cecy = int(ch_ext.get('cy', ecy)) if ch_ext is not None else ecy

    return (ox, oy, ecx, ecy, cox, coy, cecx, cecy)


def _iter_slide_shapes(container, type_filter=None):
    """Recursively yield _AbsShape proxies from container, penetrating GROUPs.

    type_filter: tuple of MSO_SHAPE_TYPE to include (None = all non-GROUP).
    Yields _AbsShape with absolute slide coordinates for left/top.
    """
    for shape in container.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            ox, oy, ecx, ecy, cox, coy, cecx, cecy = _get_group_xfrm(shape)
            for child_proxy in _iter_slide_shapes(shape, type_filter):
                # child_proxy.left/top are in this group's child coord space
                cx, cy = child_proxy.left, child_proxy.top
                abs_x = int(ox + (cx - cox) * ecx / cecx) if cecx else ox + cx
                abs_y = int(oy + (cy - coy) * ecy / cecy) if cecy else oy + cy
                raw = child_proxy._shape if isinstance(child_proxy, _AbsShape) else child_proxy
                yield _AbsShape(raw, abs_x, abs_y)
        else:
            if type_filter is None or shape.shape_type in type_filter:
                yield _AbsShape(shape, shape.left, shape.top)


def _iter_text_shapes(slide):
    """Yield all shapes with text (as _AbsShape), including inside GROUPs (recursive)."""
    for proxy in _iter_slide_shapes(slide):
        if hasattr(proxy, 'has_text_frame') and proxy.has_text_frame:
            yield proxy


def _conf_rank(confidence: str) -> int:
    """Rank confidence levels for comparison."""
    return {"high": 3, "medium": 2, "low": 1}.get(confidence, 0)


_PLAN_KEYWORDS = ("結構平面圖", "結構配置圖", "配置圖", "平面圖")


def scan_floors(prs) -> dict[int, list[dict]]:
    """Scan all slides for floor label text with confidence scoring.

    Returns {slide_num: [{"label": str, "confidence": str,
                          "source_text": str, "position": str}]} (1-based).

    Confidence:
      - "high": label in title region (top 20%) or co-occurs with plan keywords
      - "medium": label in body area without special context
      - "low": label in legend region (likely a section label, not floor title)
    """
    result = {}
    slide_h = prs.slide_height

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        detections = []
        for shape in _iter_text_shapes(slide):
            text = shape.text_frame.text.strip()
            if not text:
                continue

            # Shape position info
            shape_top_pct = shape.top / slide_h if slide_h > 0 else 0
            is_title_region = shape_top_pct < 0.20  # Top 20% of slide

            has_plan_keyword = any(kw in text for kw in _PLAN_KEYWORDS)

            # Check if shape is in legend region (heuristic: contains legend keywords)
            in_legend = any(kw in text for kw in LEGEND_KEYWORDS)

            for line in text.split("\n"):
                line_s = line.strip()
                for m in FLOOR_LABEL_RE.finditer(line_s):
                    label = m.group(1).strip()

                    # Confidence scoring
                    if is_title_region or has_plan_keyword:
                        confidence = "high"
                    elif in_legend:
                        confidence = "low"
                    else:
                        confidence = "medium"

                    pos = "title" if is_title_region else f"y={shape_top_pct:.0%}"
                    detections.append({
                        "label": label,
                        "confidence": confidence,
                        "source_text": line_s[:80],
                        "position": pos,
                    })

        if detections:
            # Deduplicate labels, keeping highest confidence for each
            seen: dict[str, dict] = {}
            for d in detections:
                lbl = d["label"]
                if lbl not in seen or _conf_rank(d["confidence"]) > _conf_rank(seen[lbl]["confidence"]):
                    seen[lbl] = d
            result[slide_num] = list(seen.values())
    return result


def scan_floors_labels(prs) -> dict[int, list[str]]:
    """Convenience wrapper: scan_floors() -> {slide_num: [label_strings]}.

    Filters out low-confidence detections (likely legend labels).
    """
    raw = scan_floors(prs)
    result = {}
    for sn, detections in raw.items():
        # Prefer medium/high confidence
        labels = [d["label"] for d in detections if d["confidence"] != "low"]
        if not labels:
            # Fallback to all if nothing passes filter
            labels = [d["label"] for d in detections]
        # Deduplicate preserving order
        seen = set()
        deduped = []
        for lbl in labels:
            if lbl not in seen:
                seen.add(lbl)
                deduped.append(lbl)
        if deduped:
            result[sn] = deduped
    return result


def format_scan_floors_output(detected: dict[int, list[dict]]) -> str:
    """Format detected floor labels (with confidence) for display.

    Args:
        detected: {slide_num: [{"label", "confidence", "source_text", "position"}]}
    """
    lines = ["Detected floor mapping:"]
    page_parts = []
    for sn in sorted(detected.keys()):
        entries = detected[sn]
        for entry in entries:
            label = entry["label"]
            conf = entry["confidence"]
            src = entry["source_text"]
            lines.append(f"  Slide {sn}: {label} ({conf} confidence, source: \"{src}\")")

        # For --page-floors suggestion, use first label
        labels = [e["label"] for e in entries]
        page_parts.append(f"{sn}={labels[0]}" if len(labels) == 1
                          else f"{sn}={labels[0]}~{labels[-1]}")

    lines.append("")
    lines.append("Suggested --page-floors:")
    lines.append(f'  "{", ".join(page_parts)}"')
    return "\n".join(lines)


def parse_legend_label(label: str):
    """Parse a legend label text.

    Returns (element_type, section_or_None, specificity, prefix, is_diaphragm).
    """
    lab = label.strip()

    # 1) Specific section: SB30x60, FB100x230, FWB60x230, B55x80 …
    m = _SECTION_RE.search(lab)
    if m:
        prefix = m.group(1).upper()
        w, d = int(m.group(2)), int(m.group(3))
        section = f"{prefix}{w}X{d}"
        if prefix in ("SB", "FSB"):
            return "small_beam", section, 1, prefix, False
        return "beam", section, 1, prefix, False

    # 2) Wall thickness: "90CM 連續壁"
    m = _WALL_CM_RE.search(lab)
    if m:
        t = int(m.group(1))
        return "wall", f"W{t}", 1, "W", True

    # 3) Generic Chinese keywords
    if "連續壁" in lab:
        return "wall", None, 0, "W", True
    if "<RC大梁>" in lab or "大梁" in lab:
        return "beam", None, 0, "B", False
    if "<RC小梁>" in lab or "小梁" in lab or "次梁" in lab:
        return "small_beam", None, 0, "SB", False
    if any(kw in lab for kw in ("<RC柱>", "邊柱", "內柱", "角柱")):
        return "column", None, 0, "C", False
    if "壁" in lab or "剪力牆" in lab:
        return "wall", None, 0, "W", False

    return "unknown", None, 0, "", False


def _round5(val_m: float) -> int:
    """Round meter value to nearest 5 cm (returns cm)."""
    cm = val_m * 100
    return max(5, round(cm / ROUND_5CM) * ROUND_5CM)


def _direction_of(x1, y1, x2, y2) -> str:
    dx, dy = abs(x2 - x1), abs(y2 - y1)
    if dx < 0.01 and dy < 0.01:
        return ""
    if dy < 0.01 or (dx > 0 and dy / max(dx, 1e-9) < DIRECTION_RATIO):
        return "X"
    if dx < 0.01 or (dy > 0 and dx / max(dy, 1e-9) < DIRECTION_RATIO):
        return "Y"
    return ""


def _wall_centerline(verts):
    """Extract (x1,y1,x2,y2) centerline and thickness from wall polygon.

    For a 4-vertex polygon: centerline = midpoints of the shorter pair of
    opposite edges; thickness = average of that shorter pair's lengths.
    For other polygons: bounding-box heuristic.

    Returns ((x1,y1,x2,y2), thickness_m) or (None, 0).
    """
    if len(verts) < 3:
        return None, 0

    if len(verts) == 4:
        def elen(i, j):
            return math.hypot(verts[j][0] - verts[i][0],
                              verts[j][1] - verts[i][1])

        e01, e12, e23, e30 = elen(0, 1), elen(1, 2), elen(2, 3), elen(3, 0)
        pair1 = (e01 + e23) / 2   # edges 0-1, 2-3
        pair2 = (e12 + e30) / 2   # edges 1-2, 3-0

        if pair1 > pair2:
            # pair2 = short edges → thickness; centerline through pair2 midpoints
            thickness = pair2
            m1 = ((verts[3][0] + verts[0][0]) / 2,
                  (verts[3][1] + verts[0][1]) / 2)
            m2 = ((verts[1][0] + verts[2][0]) / 2,
                  (verts[1][1] + verts[2][1]) / 2)
        else:
            thickness = pair1
            m1 = ((verts[0][0] + verts[1][0]) / 2,
                  (verts[0][1] + verts[1][1]) / 2)
            m2 = ((verts[2][0] + verts[3][0]) / 2,
                  (verts[2][1] + verts[3][1]) / 2)

        return ((round(m1[0], COORD_PRECISION), round(m1[1], COORD_PRECISION),
                 round(m2[0], COORD_PRECISION), round(m2[1], COORD_PRECISION)),
                thickness)

    # Fallback: bounding box
    xs = [v[0] for v in verts]
    ys = [v[1] for v in verts]
    mn_x, mx_x = min(xs), max(xs)
    mn_y, mx_y = min(ys), max(ys)
    w, h = mx_x - mn_x, mx_y - mn_y
    if w > h:
        thickness = h
        cy = (mn_y + mx_y) / 2
        return ((round(mn_x, COORD_PRECISION), round(cy, COORD_PRECISION),
                 round(mx_x, COORD_PRECISION), round(cy, COORD_PRECISION)),
                thickness)
    else:
        thickness = w
        cx = (mn_x + mx_x) / 2
        return ((round(cx, COORD_PRECISION), round(mn_y, COORD_PRECISION),
                 round(cx, COORD_PRECISION), round(mx_y, COORD_PRECISION)),
                thickness)


def _elem_key(e: dict) -> tuple:
    """Key for deduplication: type + coordinates + section."""
    return (
        e["element_type"],
        round(e["x1"], COORD_PRECISION),
        round(e["y1"], COORD_PRECISION),
        round(e["x2"], COORD_PRECISION),
        round(e["y2"], COORD_PRECISION),
        e.get("section", ""),
    )


def _snap_walls_to_beams(elements, tolerance=1.0, warnings=None):
    """Snap wall fixed-axis coordinates to nearest beam sharing the same direction.

    On 2D structural plans, walls/連續壁 are drawn parallel to but offset from beams.
    In the model they should share the beam's axis coordinate.

    Args:
        elements: list of element dicts (modified in-place).
        tolerance: max snap distance in meters (default 1.0m).
        warnings: list to append snap messages to.
    """
    beams = [e for e in elements if e["element_type"] == "beam"]
    walls = [e for e in elements if e["element_type"] == "wall"]

    snap_count = 0
    for wall in walls:
        direction = _direction_of(wall["x1"], wall["y1"], wall["x2"], wall["y2"])
        if direction == "X":
            # X-direction wall: fixed Y -> snap to nearest X-direction beam's Y
            fixed_coord = wall["y1"]
            candidates = [b for b in beams
                          if _direction_of(b["x1"], b["y1"], b["x2"], b["y2"]) == "X"]
            if not candidates:
                continue
            best = min(candidates, key=lambda b: abs(b["y1"] - fixed_coord))
            dist = abs(best["y1"] - fixed_coord)
            if 0.01 < dist <= tolerance:
                old_y = wall["y1"]
                wall["y1"] = best["y1"]
                wall["y2"] = best["y1"]  # Fixed axis: both endpoints same Y
                snap_count += 1
                msg = (f"Wall at Y={old_y:.2f} snapped to beam at "
                       f"Y={best['y1']:.2f} (delta={dist:.2f}m)")
                if warnings is not None:
                    warnings.append(msg)
        elif direction == "Y":
            # Y-direction wall: fixed X -> snap to nearest Y-direction beam's X
            fixed_coord = wall["x1"]
            candidates = [b for b in beams
                          if _direction_of(b["x1"], b["y1"], b["x2"], b["y2"]) == "Y"]
            if not candidates:
                continue
            best = min(candidates, key=lambda b: abs(b["x1"] - fixed_coord))
            dist = abs(best["x1"] - fixed_coord)
            if 0.01 < dist <= tolerance:
                old_x = wall["x1"]
                wall["x1"] = best["x1"]
                wall["x2"] = best["x1"]  # Fixed axis: both endpoints same X
                snap_count += 1
                msg = (f"Wall at X={old_x:.2f} snapped to beam at "
                       f"X={best['x1']:.2f} (delta={dist:.2f}m)")
                if warnings is not None:
                    warnings.append(msg)

    if snap_count > 0:
        print(f"  Wall snap: {snap_count} walls snapped to beam coordinates")


def group_and_dedup(all_elements: list[dict]) -> dict[str, list[dict]]:
    """Group elements by type; merge floors for duplicates (same position+section).

    Returns {"columns": [...], "beams": [...], "walls": [...], "small_beams": [...]}.
    """
    merged: dict[tuple, dict] = {}

    for e in all_elements:
        key = _elem_key(e)
        if key in merged:
            # Merge floors (union, preserve order)
            existing_floors = merged[key]["floors"]
            for f in e["floors"]:
                if f not in existing_floors:
                    existing_floors.append(f)
        else:
            merged[key] = dict(e)  # shallow copy

    # Split into typed arrays
    result = {"columns": [], "beams": [], "walls": [], "small_beams": []}
    for elem in merged.values():
        etype = elem.pop("element_type")
        # Preserve page_num for affine calibration (Phase 2 needs it)
        elem.pop("section_uncertain", None)
        target_key = {
            "column": "columns",
            "beam": "beams",
            "wall": "walls",
            "small_beam": "small_beams",
        }.get(etype)
        if target_key:
            # For columns: use grid_x/grid_y naming
            if etype == "column":
                out = {
                    "grid_x": elem["x1"],
                    "grid_y": elem["y1"],
                    "section": elem["section"],
                    "floors": elem["floors"],
                }
                if "page_num" in elem:
                    out["page_num"] = elem["page_num"]
                result["columns"].append(out)
            else:
                out = {
                    "x1": elem["x1"], "y1": elem["y1"],
                    "x2": elem["x2"], "y2": elem["y2"],
                    "section": elem["section"],
                    "floors": elem["floors"],
                    "direction": elem.get("direction", ""),
                }
                if "page_num" in elem:
                    out["page_num"] = elem["page_num"]
                if etype == "wall" and elem.get("is_diaphragm_wall"):
                    out["is_diaphragm_wall"] = True
                result[target_key].append(out)

    return result


def collect_sections(grouped: dict) -> dict:
    """Gather all unique section names by type."""
    frame_set = set()
    wall_set = set()

    for c in grouped.get("columns", []):
        if c["section"]:
            frame_set.add(c["section"])
    for b in grouped.get("beams", []):
        if b["section"]:
            frame_set.add(b["section"])
    for sb in grouped.get("small_beams", []):
        if sb["section"]:
            frame_set.add(sb["section"])
    for w in grouped.get("walls", []):
        s = w["section"]
        if s:
            m = re.match(r"^W(\d+)$", s)
            if m:
                wall_set.add(int(m.group(1)))

    return {
        "frame": sorted(frame_set),
        "wall": sorted(wall_set),
    }


# ─── Constants ────────────────────────────────────────────────────────────────

_COLOR_TOLERANCE_OVERRIDE = None   # Set by --color-tolerance CLI; None = use default (15)
TICK_MAX_LENGTH_EMU = 300000       # Max length for a tick mark (short line)
TICK_Y_TOLERANCE_EMU = 100000     # Y-proximity for grouping tick marks
TEXT_TICK_Y_TOLERANCE_EMU = 300000 # Y-proximity for matching text to ticks
COLUMN_DEDUP_TOLERANCE_EMU = 5000 # Position tolerance for column dedup
LEGEND_KEYWORDS = [
    "RC", "大梁", "小梁", "柱", "壁", "連續壁", "剪力牆",
    "梁", "版", "板", "基礎", "FB", "WB", "SB",
]
# Phase 2 specific keywords for fallback legend scan
LEGEND_KEYWORDS_PHASE2 = ["SB", "FSB", "小梁", "次梁", "版", "板", "SLAB", "FS"]
# Phase 2 label regex for whole-slide scan fallback
_PHASE2_LABEL_RE = re.compile(
    r"(FSB|SB)\d+\s*[xX]\s*\d+|FS\d+|S\d+|小梁|次梁|版|板|SLAB",
    re.IGNORECASE,
)
SLIDE_AREA_RATIO_MAX = 0.50       # Exclude shapes > 50% of slide area
MIN_SCALE_MEASUREMENTS = 1        # Require at least 1 measurement per page
MIN_COLUMN_DIM_EMU = 30000        # Min column shape dimension (~3cm at typical scale)
COLUMN_REQUIRED_VERTICES = (4, 5)  # 4- or 5-vertex freeforms (5 = explicit close point)
MEASUREMENT_RE = re.compile(r"(\d+\.?\d*)\s*m\b")
# Column section pattern: C100x100, C110x140, etc.
_COLUMN_SECTION_RE = re.compile(r"C(\d+)\s*[xX]\s*(\d+)")


# ─── Freeform Vertex Counting ────────────────────────────────────────────────

_VERTEX_NS = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}


def _get_freeform_vertex_count(shape):
    """Count vertices in a FREEFORM shape's custom geometry path.

    Returns vertex count (moveTo + lnTo nodes, excluding close),
    or None if geometry data is not accessible.
    """
    try:
        el = shape._element
    except AttributeError:
        return None

    paths = el.findall('.//a:custGeom/a:pathLst/a:path', _VERTEX_NS)
    if not paths:
        return None

    count = 0
    for path in paths:
        count += len(path.findall('a:moveTo', _VERTEX_NS))
        count += len(path.findall('a:lnTo', _VERTEX_NS))
    return count


# ─── Theme Color Resolution ──────────────────────────────────────────────────

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# Map theme XML child tag names to MSO_THEME_COLOR_INDEX enum values
_THEME_TAG_TO_ENUM = {
    "dk1": MSO_THEME_COLOR_INDEX.DARK_1,
    "lt1": MSO_THEME_COLOR_INDEX.LIGHT_1,
    "dk2": MSO_THEME_COLOR_INDEX.DARK_2,
    "lt2": MSO_THEME_COLOR_INDEX.LIGHT_2,
    "accent1": MSO_THEME_COLOR_INDEX.ACCENT_1,
    "accent2": MSO_THEME_COLOR_INDEX.ACCENT_2,
    "accent3": MSO_THEME_COLOR_INDEX.ACCENT_3,
    "accent4": MSO_THEME_COLOR_INDEX.ACCENT_4,
    "accent5": MSO_THEME_COLOR_INDEX.ACCENT_5,
    "accent6": MSO_THEME_COLOR_INDEX.ACCENT_6,
    "hlink": MSO_THEME_COLOR_INDEX.HYPERLINK,
    "folHlink": MSO_THEME_COLOR_INDEX.FOLLOWED_HYPERLINK,
}


def _build_theme_color_map(prs):
    """Parse PPT theme XML to build theme_color enum → RGB hex mapping.

    Returns dict: {MSO_THEME_COLOR_INDEX.ACCENT_1: "4F81BD", ...}
    """
    try:
        sm = prs.slide_masters[0]
        theme_part = None
        for rel in sm.part.rels.values():
            if "theme" in rel.reltype:
                theme_part = rel.target_part
                break
        if theme_part is None:
            return {}

        root = etree.fromstring(theme_part.blob)
        clr_scheme = root.find(f".//{{{_A_NS}}}clrScheme")
        if clr_scheme is None:
            return {}

        result = {}
        for child in clr_scheme:
            tag_local = etree.QName(child.tag).localname
            enum_val = _THEME_TAG_TO_ENUM.get(tag_local)
            if enum_val is None:
                continue
            # Read srgbClr or sysClr
            srgb = child.find(f"{{{_A_NS}}}srgbClr")
            if srgb is not None:
                result[enum_val] = srgb.get("val", "").upper()
            else:
                sys_clr = child.find(f"{{{_A_NS}}}sysClr")
                if sys_clr is not None:
                    result[enum_val] = sys_clr.get("lastClr", "000000").upper()
        return result
    except Exception:
        return {}


def _apply_brightness(hex_color, brightness):
    """Apply brightness adjustment to a hex color string.

    brightness > 0: tint (blend toward white)
    brightness < 0: shade (blend toward black)
    """
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    if brightness > 0:
        # Tint: blend toward 255
        r = int(r + (255 - r) * brightness)
        g = int(g + (255 - g) * brightness)
        b = int(b + (255 - b) * brightness)
    elif brightness < 0:
        # Shade: blend toward 0
        factor = 1 + brightness  # brightness is negative
        r = int(r * factor)
        g = int(g * factor)
        b = int(b * factor)
    r = max(0, min(255, r))
    g = max(0, min(255, g))
    b = max(0, min(255, b))
    return f"{r:02X}{g:02X}{b:02X}"


def _resolve_scheme_color(color_obj, theme_map):
    """Resolve a python-pptx color object with type==SCHEME to RGB hex string."""
    try:
        tc = color_obj.theme_color
        hex_val = theme_map.get(tc)
        if hex_val is None:
            return None
        brightness = color_obj.brightness
        if brightness and brightness != 0:
            hex_val = _apply_brightness(hex_val, brightness)
        return hex_val
    except Exception:
        return None


# ─── Scale Detection ─────────────────────────────────────────────────────────

def _get_shape_color(shape, color_type="line", theme_map=None):
    """Extract RGB hex string from a shape's line or fill color.

    Handles RGB colors directly and resolves SCHEME (theme) colors
    via theme_map when provided. SCHEME colors without theme_map are
    skipped (python-pptx raises on .rgb for scheme colors).
    """
    try:
        if color_type == "line":
            ln = shape.line
            if ln and ln.color and ln.color.type is not None:
                if ln.color.type == MSO_COLOR_TYPE.SCHEME:
                    if theme_map:
                        return _resolve_scheme_color(ln.color, theme_map)
                    return None  # Can't resolve scheme without theme_map
                if ln.color.rgb:
                    return str(ln.color.rgb)
        elif color_type == "fill":
            fill = shape.fill
            if fill.type is not None:
                fc = fill.fore_color
                if fc and fc.type is not None:
                    if fc.type == MSO_COLOR_TYPE.SCHEME:
                        if theme_map:
                            return _resolve_scheme_color(fc, theme_map)
                        return None  # Can't resolve scheme without theme_map
                    if fc.rgb:
                        return str(fc.rgb)
    except Exception:
        pass
    return None


def _find_tick_marks(slide):
    """Find short line segments (tick marks) on a slide.

    Returns list of (orientation, x, y, length, color) where:
      orientation: 'V' (vertical, width=0) or 'H' (horizontal, height=0)
      x, y: left/top position in EMU
      length: in EMU
    """
    ticks = []
    for shape in _iter_slide_shapes(slide, (MSO_SHAPE_TYPE.FREEFORM,)):
        w, h = shape.width, shape.height
        if w == 0 and 0 < h <= TICK_MAX_LENGTH_EMU:
            color = _get_shape_color(shape, "line")
            ticks.append(("V", shape.left, shape.top, h, color))
        elif h == 0 and 0 < w <= TICK_MAX_LENGTH_EMU:
            color = _get_shape_color(shape, "line")
            ticks.append(("H", shape.left, shape.top, w, color))
    return ticks


def _find_measurement_texts(slide):
    """Find TEXT_BOX shapes containing measurement values like '8.5 m'.

    Returns list of (meters_value, center_x, center_y, text).
    """
    measurements = []
    for shape in _iter_slide_shapes(slide, (MSO_SHAPE_TYPE.TEXT_BOX,)):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        # Only match simple single-value measurement texts
        # Skip multi-line texts like "27.4 m\n9.0 m\t8.6 m\t9.8 m"
        lines = text.split("\n")
        for line in lines:
            # Skip lines with tabs (multi-value)
            if "\t" in line:
                continue
            m = MEASUREMENT_RE.match(line.strip())
            if m:
                val = float(m.group(1))
                if val < 1.0 or val > 100.0:
                    continue  # Skip unreasonable values
                cx = shape.left + shape.width // 2
                cy = shape.top + shape.height // 2
                measurements.append((val, cx, cy, line.strip()))
    return measurements


def _estimate_scale_from_shapes(slide, slide_num):
    """Estimate rough EMU/m scale from FREEFORM shape bounding box.

    Fallback when no measurement texts are available.
    Assumes the drawing extent represents ~30m (typical structural plan).
    """
    min_x, min_y = float('inf'), float('inf')
    max_x, max_y = float('-inf'), float('-inf')
    count = 0
    for shape in _iter_slide_shapes(slide, (MSO_SHAPE_TYPE.FREEFORM,)):
        left, top = shape.left, shape.top
        right = left + shape.width
        bottom = top + shape.height
        min_x = min(min_x, left)
        min_y = min(min_y, top)
        max_x = max(max_x, right)
        max_y = max(max_y, bottom)
        count += 1

    if count < 5:
        return None, f"Slide {slide_num}: too few shapes ({count}) for scale estimation"

    extent_x = max_x - min_x
    extent_y = max_y - min_y
    extent = max(extent_x, extent_y)

    if extent < 1000000:  # Less than ~1m worth of EMU
        return None, f"Slide {slide_num}: shape extent too small for estimation"

    # Assume the longer drawing dimension ≈ 30m
    ASSUMED_SPAN_M = 30.0
    scale = extent / ASSUMED_SPAN_M

    return scale, {
        "slide": slide_num,
        "emu_per_meter": round(scale, 1),
        "num_measurements": 0,
        "estimated": True,
        "assumed_span_m": ASSUMED_SPAN_M,
        "extent_emu": round(extent, 0),
    }


def detect_slide_scale(slide, slide_num):
    """Detect the EMU-per-meter scale for a slide.

    Strategy:
    1. Find measurement texts (e.g., "8.5 m", "11.0 m")
    2. Find vertical tick marks (short vertical lines near measurements)
    3. Group V-ticks by similar Y into dimension lines
    4. For each measurement text, find the best matching V-tick pair
    5. Return the average scale (EMU per meter)

    Returns (emu_per_meter, details_dict) or (None, error_msg).
    """
    measurements = _find_measurement_texts(slide)
    if not measurements:
        return None, f"Slide {slide_num}: no measurement texts found (need 'X.X m' TEXT_BOX)"

    ticks = _find_tick_marks(slide)
    v_ticks = [(x, y, length) for orient, x, y, length, _ in ticks if orient == "V"]

    if len(v_ticks) < 2:
        return None, f"Slide {slide_num}: fewer than 2 vertical tick marks found"

    # Sort V-ticks by X position
    v_ticks.sort(key=lambda t: t[0])

    # Try to match each measurement text to a pair of V-ticks
    scales = []
    details = []

    for meters, text_cx, text_cy, text_str in measurements:
        best_pair = None
        best_score = float("inf")

        for i in range(len(v_ticks)):
            for j in range(i + 1, len(v_ticks)):
                x_i, y_i, _ = v_ticks[i]
                x_j, y_j, _ = v_ticks[j]

                # V-ticks should be at similar Y
                if abs(y_i - y_j) > TICK_Y_TOLERANCE_EMU:
                    continue

                # Text should be between the ticks (X-wise) or close
                tick_cx = (x_i + x_j) / 2
                tick_cy = (y_i + y_j) / 2

                # Text should be near the tick Y
                if abs(text_cy - tick_cy) > TEXT_TICK_Y_TOLERANCE_EMU:
                    continue

                # Text center should be between ticks (with margin)
                margin = (x_j - x_i) * 0.3
                if text_cx < x_i - margin or text_cx > x_j + margin:
                    continue

                # Score: prefer pair where text is centered between ticks
                dx = abs(text_cx - tick_cx)
                dy = abs(text_cy - tick_cy)
                score = dx + dy

                if score < best_score:
                    best_score = score
                    best_pair = (x_i, x_j, meters)

        if best_pair:
            x_left, x_right, m = best_pair
            emu_dist = x_right - x_left
            scale = emu_dist / m
            scales.append(scale)
            details.append({
                "text": text_str,
                "meters": m,
                "emu_distance": emu_dist,
                "emu_per_meter": round(scale, 1),
            })

    if not scales:
        return None, f"Slide {slide_num}: could not pair any measurement text with tick marks"

    # Use median scale for robustness
    scales.sort()
    median_scale = scales[len(scales) // 2]

    return median_scale, {
        "slide": slide_num,
        "measurements": details,
        "emu_per_meter": round(median_scale, 1),
        "num_measurements": len(scales),
    }


# ─── Coordinate Conversion ───────────────────────────────────────────────────

def _emu_to_m(x_emu, y_emu, scale, slide_h):
    """Convert EMU coordinates to meters.

    PPT origin is top-left, Y increases downward.
    Structural origin is bottom-left, Y increases upward.
    """
    x_m = round(x_emu / scale, COORD_PRECISION)
    y_m = round((slide_h - y_emu) / scale, COORD_PRECISION)
    return x_m, y_m


# ─── PNG Extraction ──────────────────────────────────────────────────────────

def extract_base_png(slide, slide_num, floor_label, crop_dir):
    """Extract PNG images from GROUP shapes on the slide.

    Saves as {floor_label}_full.png in crop_dir.
    Returns list of saved file paths.
    """
    saved = []
    crop_path = Path(crop_dir)
    crop_path.mkdir(parents=True, exist_ok=True)

    pic_idx = 0
    for shape in slide.shapes:
        # Extract from GROUP children
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img = child.image
                    ext = {
                        "image/png": ".png",
                        "image/jpeg": ".jpg",
                        "image/gif": ".gif",
                    }.get(img.content_type, ".png")
                    fname = f"{floor_label}_full{ext}" if pic_idx == 0 else f"{floor_label}_full_{pic_idx}{ext}"
                    fpath = crop_path / fname
                    with open(fpath, "wb") as f:
                        f.write(img.blob)
                    saved.append(str(fpath))
                    print(f"  Saved: {fpath} ({len(img.blob)} bytes)")
                    pic_idx += 1
        # Extract top-level PICTURE shapes
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img = shape.image
            ext = {
                "image/png": ".png",
                "image/jpeg": ".jpg",
                "image/gif": ".gif",
            }.get(img.content_type, ".png")
            fname = f"{floor_label}_full{ext}" if pic_idx == 0 else f"{floor_label}_full_{pic_idx}{ext}"
            fpath = crop_path / fname
            with open(fpath, "wb") as f:
                f.write(img.blob)
            saved.append(str(fpath))
            print(f"  Saved: {fpath} ({len(img.blob)} bytes)")
            pic_idx += 1
    return saved


# ─── Legend Parsing ──────────────────────────────────────────────────────────

def _detect_legend_region(slide, slide_w, slide_h):
    """Detect whether the legend is on the left or right side of the slide.

    Strategy: Count structural keywords in text shapes on each side.
    The side with more keywords is the legend side.

    Returns ("left", x_boundary) or ("right", x_boundary).
    """
    mid_x = slide_w // 2
    left_score = 0
    right_score = 0

    for shape in _iter_text_shapes(slide):
        text = shape.text_frame.text
        cx = shape.left + shape.width // 2

        score = sum(1 for kw in LEGEND_KEYWORDS if kw in text)
        if score > 0:
            if cx < mid_x * 0.6:  # Clearly on left side
                left_score += score
            elif cx > mid_x * 1.4:  # Clearly on right side
                right_score += score

    # Default to left if can't determine
    if left_score >= right_score:
        # Legend on left: find the rightmost legend text to set boundary
        max_right = 0
        for shape in _iter_text_shapes(slide):
            text = shape.text_frame.text
            if any(kw in text for kw in LEGEND_KEYWORDS):
                r = shape.left + shape.width
                if r > max_right and shape.left < mid_x:
                    max_right = r
        boundary = max_right + 200000 if max_right > 0 else int(slide_w * 0.25)
        return "left", boundary
    else:
        min_left = slide_w
        for shape in _iter_text_shapes(slide):
            text = shape.text_frame.text
            if any(kw in text for kw in LEGEND_KEYWORDS):
                if shape.left > mid_x:
                    min_left = min(min_left, shape.left)
        boundary = min_left - 200000 if min_left < slide_w else int(slide_w * 0.75)
        return "right", boundary


def _match_label_to_swatch(lbl, swatches, dy_tol=300000, dx_tol=500000):
    """Find nearest swatch for a label entry within tolerance."""
    best_swatch = None
    best_dist = float("inf")
    for sw in swatches:
        dy = abs(sw["cy"] - lbl["cy"])
        if dy > dy_tol:
            continue
        dx = sw["cx"] - lbl["cx"]
        if dx > dx_tol:
            continue
        dist = math.sqrt(dx ** 2 + dy ** 2)
        if dist < best_dist:
            best_dist = dist
            best_swatch = sw
    return best_swatch


def _parse_text_to_labels(legend_texts):
    """Parse text blocks into label entries (shared by Pass 1 and Pass 2)."""
    label_entries = []
    for lt in legend_texts:
        text = lt["text"]
        lines = text.split("\n")
        y_offset = 0
        line_height = (lt["bottom"] - lt["top"]) / max(len(lines), 1)
        for line in lines:
            line = line.strip()
            if not line:
                y_offset += 1
                continue
            etype, section, spec, prefix, is_diaphragm = parse_legend_label(line)
            # PPT-specific: handle column section labels like "C100x100 C100x120"
            if etype == "unknown":
                cm = _COLUMN_SECTION_RE.search(line)
                if cm:
                    w_val, d_val = int(cm.group(1)), int(cm.group(2))
                    lo, hi = sorted([w_val, d_val])
                    etype = "column"
                    section = f"C{lo}X{hi}"
                    spec = 1
                    prefix = "C"
                    is_diaphragm = False
            if etype != "unknown":
                label_entries.append({
                    "label": line,
                    "element_type": etype,
                    "section": section,
                    "specificity": spec,
                    "prefix": prefix,
                    "is_diaphragm": is_diaphragm,
                    "cx": lt["cx"],
                    "cy": int(lt["top"] + y_offset * line_height + line_height / 2),
                    "left": lt["left"],
                })
            y_offset += 1
    return label_entries


def _collect_legend_swatches(slide, slide_w, slide_h, side, boundary,
                             theme_map=None):
    """Collect colored swatches in legend region."""
    swatches = []
    for shape in _iter_slide_shapes(slide, (MSO_SHAPE_TYPE.FREEFORM,
                                            MSO_SHAPE_TYPE.AUTO_SHAPE)):
        cx = shape.left + shape.width // 2
        if side == "left" and cx > boundary:
            continue
        if side == "right" and cx < boundary:
            continue
        if shape.width > slide_w * 0.5 or shape.height > slide_h * 0.5:
            continue
        line_color = _get_shape_color(shape, "line", theme_map)
        fill_color = _get_shape_color(shape, "fill", theme_map)
        # Prefer fill color (represents element's actual color);
        # outline is usually just decorative (e.g. black border)
        color = fill_color or line_color
        if color and color != "FFFFFF":
            swatches.append({
                "color": color,
                "fill_color": fill_color,
                "line_color": line_color,
                "left": shape.left,
                "top": shape.top,
                "cx": cx,
                "cy": shape.top + shape.height // 2,
                "right": shape.left + shape.width,
            })
    return swatches


def _get_cell_bg_color(cell, theme_map=None):
    """Extract background color hex from a table cell.

    Tries cell.fill.fore_color (RGB or SCHEME).
    Returns uppercase hex string like "FF0000", or None.
    """
    try:
        fill = cell.fill
        if fill.type is None:
            return None
        fc = fill.fore_color
        if fc is None or fc.type is None:
            return None
        if fc.type == MSO_COLOR_TYPE.SCHEME:
            if theme_map:
                return _resolve_scheme_color(fc, theme_map)
            return None  # Can't resolve scheme without theme_map
        if fc.rgb:
            return str(fc.rgb)
    except Exception:
        pass
    return None


def _parse_table_legend(slide, slide_w, phase, theme_map=None):
    """Parse legend from PPT table (2-col: [bg_color | section_name]).

    Returns (legend_dict, diagnostics, boundary, side) or None if no table found.
    """
    table_shape = None
    for shape in _iter_slide_shapes(slide):
        if hasattr(shape, 'has_table') and shape.has_table and len(shape.table.columns) == 2:
            table_shape = shape
            break

    if table_shape is None:
        return None

    table = table_shape.table
    mapping = {}
    entries_diag = []

    for row in table.rows:
        cell_0 = row.cells[0]
        cell_1 = row.cells[1]

        # Get background color from first cell
        color_hex = _get_cell_bg_color(cell_0, theme_map)
        if not color_hex or color_hex == "FFFFFF":
            continue

        # Get section text from second cell
        text = cell_1.text.strip()
        if not text:
            continue

        # Parse using existing label parser
        etype, section, spec, prefix, is_diaphragm = parse_legend_label(text)

        # PPT-specific: handle column section labels
        if etype == "unknown":
            cm = _COLUMN_SECTION_RE.search(text)
            if cm:
                w_val, d_val = int(cm.group(1)), int(cm.group(2))
                lo, hi = sorted([w_val, d_val])
                etype = "column"
                section = f"C{lo}X{hi}"
                spec = 1
                prefix = "C"
                is_diaphragm = False

        if etype == "unknown":
            continue

        # Phase filter
        if phase == "phase1" and etype == "small_beam":
            continue
        if phase == "phase2" and etype not in ("small_beam",):
            continue

        rgb = [int(color_hex[i:i+2], 16) for i in (0, 2, 4)]
        entry = LegendEntry(
            element_type=etype,
            section=section,
            color_name=color_hex,
            color_rgb=rgb,
            specificity=spec,
            is_diaphragm=is_diaphragm,
            prefix=prefix,
            label=text,
        )
        mapping.setdefault(color_hex, []).append(entry)
        entries_diag.append({
            "label": text,
            "color": color_hex,
            "type": etype,
        })

    if not mapping:
        return None

    # Sort by specificity descending within each color
    for c in mapping:
        mapping[c].sort(key=lambda e: -e.specificity)

    # Determine legend side from table position
    table_cx = table_shape.left + table_shape.width // 2
    slide_cx = slide_w // 2
    if table_cx < slide_cx:
        side = "left"
        boundary = table_shape.left + table_shape.width + 200000
    else:
        side = "right"
        boundary = table_shape.left - 200000

    diagnostics = {
        "legend_source": "table",
        "table_rows": len(table.rows),
        "pass1_labels": entries_diag,
        "pass2_labels": [],
        "pass3_orphans": [],
        "all_swatches": [],
        "unmatched_labels": [],
        "unmatched_swatches": [],
    }

    return mapping, diagnostics, boundary, side


def parse_legend(slide, slide_w, slide_h, phase="all", theme_map=None):
    """Parse the legend area to build color -> LegendEntry mapping.

    Table-first strategy:
      Try PPT table legend (2-col) first; fall back to shape-based if no table.

    Shape-based three-pass strategy:
      Pass 1: Standard legend region parsing (keyword-based region + label-swatch match)
      Pass 2: (Phase 2 only) Scan ALL text shapes for SB/FSB/版/板 labels,
              retry unmatched with relaxed tolerances (dy: 500K, RGB: +/-15)
      Pass 3: Orphan swatch retry with relaxed RGB tolerance (+/-15)

    Returns (mapping, diagnostics, boundary, side).
    """
    # ── Try table-based legend first ──
    table_result = _parse_table_legend(slide, slide_w, phase, theme_map)
    if table_result is not None:
        return table_result

    # ── Existing shape-based legend parsing ──
    side, boundary = _detect_legend_region(slide, slide_w, slide_h)
    diagnostics = {
        "pass1_labels": [],
        "pass2_labels": [],
        "pass3_orphans": [],
        "all_swatches": [],
        "unmatched_labels": [],
        "unmatched_swatches": [],
    }

    # ── Collect legend text shapes in legend region ──
    legend_texts = []
    for shape in _iter_text_shapes(slide):
        cx = shape.left + shape.width // 2
        if side == "left" and cx > boundary:
            continue
        if side == "right" and cx < boundary:
            continue
        text = shape.text_frame.text.strip()
        if text:
            legend_texts.append({
                "text": text,
                "left": shape.left,
                "top": shape.top,
                "right": shape.left + shape.width,
                "bottom": shape.top + shape.height,
                "cx": cx,
                "cy": shape.top + shape.height // 2,
            })

    if not legend_texts:
        diagnostics["legend_source"] = "shape"
        return {}, diagnostics, boundary, side

    # ── Parse labels ──
    label_entries = _parse_text_to_labels(legend_texts)

    # ── Collect swatches ──
    legend_swatches = _collect_legend_swatches(slide, slide_w, slide_h, side, boundary,
                                               theme_map)
    diagnostics["all_swatches"] = [
        {"color": sw["color"], "cx": sw["cx"], "cy": sw["cy"]}
        for sw in legend_swatches
    ]

    # ════════════════════════════════════════════════════════════════
    # Pass 1: Standard matching (dy: 300K EMU, RGB: exact match on swatch)
    # ════════════════════════════════════════════════════════════════
    mapping: dict[str, list[LegendEntry]] = {}
    matched_swatch_indices = set()
    unmatched_labels_pass1 = []

    for lbl in label_entries:
        best_swatch = _match_label_to_swatch(lbl, legend_swatches)
        if best_swatch:
            idx = next(i for i, sw in enumerate(legend_swatches)
                       if sw is best_swatch)
            matched_swatch_indices.add(idx)
            color_hex = best_swatch["color"]
            rgb = [int(color_hex[i:i+2], 16) for i in (0, 2, 4)]
            entry = LegendEntry(
                element_type=lbl["element_type"],
                section=lbl["section"],
                color_name=color_hex,
                color_rgb=rgb,
                specificity=lbl["specificity"],
                is_diaphragm=lbl["is_diaphragm"],
                prefix=lbl["prefix"],
                label=lbl["label"],
            )
            mapping.setdefault(color_hex, []).append(entry)
            # Dual-color protection: if swatch has both fill and line colors
            # that differ, register the alternate color too (same entry)
            alt_color = best_swatch.get("line_color") if color_hex == best_swatch.get("fill_color") \
                else best_swatch.get("fill_color")
            if alt_color and alt_color != color_hex and alt_color != "FFFFFF" and alt_color != "000000":
                mapping.setdefault(alt_color, []).append(entry)
            diagnostics["pass1_labels"].append({
                "label": lbl["label"],
                "color": color_hex,
                "type": lbl["element_type"],
            })
        else:
            unmatched_labels_pass1.append(lbl)

    # ════════════════════════════════════════════════════════════════
    # Pass 2: Phase 2 fallback — scan ALL text shapes for SB/FSB labels
    # ════════════════════════════════════════════════════════════════
    if phase == "phase2" and unmatched_labels_pass1:
        # Scan entire slide for text shapes with phase2-relevant labels
        all_slide_texts = []
        for shape in _iter_text_shapes(slide):
            text = shape.text_frame.text.strip()
            if not text:
                continue
            # Check if text contains phase2 keywords
            if not _PHASE2_LABEL_RE.search(text):
                continue
            cx = shape.left + shape.width // 2
            all_slide_texts.append({
                "text": text,
                "left": shape.left,
                "top": shape.top,
                "right": shape.left + shape.width,
                "bottom": shape.top + shape.height,
                "cx": cx,
                "cy": shape.top + shape.height // 2,
            })

        extra_labels = _parse_text_to_labels(all_slide_texts)
        # Filter to only phase2-relevant types not already matched
        matched_sections = {e.section for entries in mapping.values()
                           for e in entries if e.section}
        extra_labels = [
            lbl for lbl in extra_labels
            if lbl["element_type"] in ("small_beam",)
            and lbl.get("section") and lbl["section"] not in matched_sections
        ]

        # Retry with relaxed tolerances (dy: 500K, RGB: ±15 via fuzzy color)
        orphan_swatches = [sw for i, sw in enumerate(legend_swatches)
                           if i not in matched_swatch_indices]

        for lbl in unmatched_labels_pass1 + extra_labels:
            best_swatch = _match_label_to_swatch(
                lbl, orphan_swatches, dy_tol=500000, dx_tol=700000)
            if best_swatch:
                color_hex = best_swatch["color"]
                # Try fuzzy RGB match: check if any existing mapping color is close
                matched_color = color_hex
                for existing_color in mapping:
                    er = int(existing_color[0:2], 16)
                    eg = int(existing_color[2:4], 16)
                    eb = int(existing_color[4:6], 16)
                    sr = int(color_hex[0:2], 16)
                    sg = int(color_hex[2:4], 16)
                    sb_ = int(color_hex[4:6], 16)
                    if max(abs(er - sr), abs(eg - sg), abs(eb - sb_)) <= 15:
                        matched_color = existing_color
                        break

                rgb = [int(matched_color[i:i+2], 16) for i in (0, 2, 4)]
                entry = LegendEntry(
                    element_type=lbl["element_type"],
                    section=lbl["section"],
                    color_name=matched_color,
                    color_rgb=rgb,
                    specificity=lbl["specificity"],
                    is_diaphragm=lbl["is_diaphragm"],
                    prefix=lbl["prefix"],
                    label=lbl["label"] + " (relaxed_match)",
                )
                mapping.setdefault(matched_color, []).append(entry)
                diagnostics["pass2_labels"].append({
                    "label": lbl["label"],
                    "color": matched_color,
                    "type": lbl["element_type"],
                    "relaxed_match": True,
                })
                # Mark swatch as used
                idx = next((i for i, sw in enumerate(legend_swatches)
                            if sw is best_swatch), None)
                if idx is not None:
                    matched_swatch_indices.add(idx)

    # ════════════════════════════════════════════════════════════════
    # Pass 3: Orphan swatch retry — fuzzy RGB (±15) for still-unmatched swatches
    # ════════════════════════════════════════════════════════════════
    final_orphan_swatches = [
        (i, sw) for i, sw in enumerate(legend_swatches)
        if i not in matched_swatch_indices
    ]

    for i, sw in final_orphan_swatches:
        sw_r = int(sw["color"][0:2], 16)
        sw_g = int(sw["color"][2:4], 16)
        sw_b = int(sw["color"][4:6], 16)

        # Try matching to an existing legend color with ±15 tolerance
        for color_hex in mapping:
            cr = int(color_hex[0:2], 16)
            cg = int(color_hex[2:4], 16)
            cb = int(color_hex[4:6], 16)
            if max(abs(sw_r - cr), abs(sw_g - cg), abs(sw_b - cb)) <= 15:
                # This orphan swatch is a fuzzy duplicate of an existing color
                diagnostics["pass3_orphans"].append({
                    "swatch_color": sw["color"],
                    "matched_to": color_hex,
                    "relaxed_match": True,
                })
                matched_swatch_indices.add(i)
                break

    # Record truly unmatched items for diagnostics
    diagnostics["unmatched_labels"] = [
        {"label": lbl["label"], "type": lbl["element_type"]}
        for lbl in unmatched_labels_pass1
        if not any(
            d["label"] == lbl["label"] for d in diagnostics["pass2_labels"]
        )
    ]
    diagnostics["unmatched_swatches"] = [
        {"color": legend_swatches[i]["color"]}
        for i in range(len(legend_swatches))
        if i not in matched_swatch_indices
    ]

    # Sort by specificity descending within each color
    for c in mapping:
        mapping[c].sort(key=lambda e: -e.specificity)

    diagnostics["legend_source"] = "shape"
    return mapping, diagnostics, boundary, side


# ─── Legend Validation ───────────────────────────────────────────────────────

def _validate_legend(slide, slide_num, legend, slide_w, slide_h,
                     legend_side, legend_boundary, warnings, theme_map=None,
                     tolerance=15):
    """Validate legend coverage: count matched/unmatched shapes per color.

    tolerance: fuzzy validation tolerance (Chebyshev distance).
               Use higher values for table-sourced legends.
    Returns validation dict for metadata.
    """
    # Count drawing-area shapes by color
    shape_colors = {}  # color_hex -> count
    for shape in _iter_slide_shapes(slide, (MSO_SHAPE_TYPE.FREEFORM,)):
        cx = shape.left + shape.width // 2
        # Skip legend region shapes
        if legend_side == "left" and cx < legend_boundary:
            continue
        if legend_side == "right" and cx > legend_boundary:
            continue
        # Skip very large shapes
        if shape.width > 0 and shape.height > 0:
            if (shape.width * shape.height) > slide_w * slide_h * SLIDE_AREA_RATIO_MAX:
                continue

        line_color = _get_shape_color(shape, "line", theme_map)
        fill_color = _get_shape_color(shape, "fill", theme_map)
        color = line_color or fill_color
        if color and color != "FFFFFF":
            shape_colors[color] = shape_colors.get(color, 0) + 1

    validation = {"slide": slide_num, "entries": []}

    # Check each legend entry
    validate_tolerance = tolerance if tolerance > 15 else 5
    for color_hex, entries in legend.items():
        count = shape_colors.pop(color_hex, 0)
        # Also check fuzzy matches
        to_remove = []
        for sc in shape_colors:
            r, g, b = int(sc[0:2], 16), int(sc[2:4], 16), int(sc[4:6], 16)
            lr, lg, lb = int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16)
            if max(abs(r-lr), abs(g-lg), abs(b-lb)) <= validate_tolerance:
                count += shape_colors[sc]
                to_remove.append(sc)
        for sc in to_remove:
            shape_colors.pop(sc, None)

        for entry in entries:
            status = "matched" if count > 0 else "no_shapes"
            validation["entries"].append({
                "color": color_hex,
                "section": entry.section or "(generic)",
                "type": entry.element_type,
                "label": entry.label,
                "shape_count": count,
                "status": status,
            })
            if count == 0:
                warnings.append(
                    f"Slide {slide_num}: legend entry #{color_hex} -> "
                    f"{entry.section or entry.label} ({entry.element_type}): "
                    f"0 shapes matched"
                )

    # Orphan shapes (drawing-area colors not in any legend)
    for color_hex, count in shape_colors.items():
        validation["entries"].append({
            "color": color_hex,
            "section": None,
            "type": "unmatched",
            "label": "",
            "shape_count": count,
            "status": "orphan",
        })
        warnings.append(
            f"Slide {slide_num}: {count} shapes with color #{color_hex} "
            f"have no legend entry"
        )

    return validation


def _print_legend_validation(validation):
    """Print legend validation report for a slide."""
    print(f"  Legend Validation (Slide {validation['slide']}):")
    for entry in validation["entries"]:
        color = entry["color"]
        section = entry["section"] or "(no legend)"
        etype = entry["type"]
        count = entry["shape_count"]
        if entry["status"] == "matched":
            print(f"    + #{color} -> {section} ({etype}): {count} shapes matched")
        elif entry["status"] == "no_shapes":
            print(f"    ? #{color} -> {section} ({etype}): 0 shapes matched")
        elif entry["status"] == "orphan":
            print(f"    ! #{color} -> no legend entry: {count} shapes unmatched")


# ─── Shape Extraction & Classification ────────────────────────────────────────

def _fuzzy_color_match(color_hex, legend, tolerance=15):
    """Find legend entry list with closest color within tolerance.

    Uses Chebyshev distance (max per-channel difference) with Manhattan
    tie-breaking when two legend colors have the same Chebyshev distance.
    Default tolerance=15 to handle PPT color variations (manual color picks,
    theme adjustments, export artifacts). Pass 3 orphan retry already used ±15;
    this unifies the main matching path to the same value.
    Returns the list of LegendEntry for the matched color, or None.
    """
    if _COLOR_TOLERANCE_OVERRIDE is not None:
        tolerance = _COLOR_TOLERANCE_OVERRIDE
    # Exact match first
    if color_hex in legend:
        return legend[color_hex]
    # Fuzzy match
    r, g, b = int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16)
    best_key, best_dist, best_manhattan = None, float('inf'), float('inf')
    for key in legend:
        kr = int(key[0:2], 16)
        kg = int(key[2:4], 16)
        kb = int(key[4:6], 16)
        dist = max(abs(r - kr), abs(g - kg), abs(b - kb))
        manhattan = abs(r - kr) + abs(g - kg) + abs(b - kb)
        if dist <= tolerance and (dist < best_dist or
                                   (dist == best_dist and manhattan < best_manhattan)):
            best_dist = dist
            best_manhattan = manhattan
            best_key = key
    return legend[best_key] if best_key else None


def _resolve_pptx_legend(color_hex, geom_type, legend, tolerance=15):
    """Resolve a shape's color to a legend entry.

    geom_type: 'line' | 'rectangle'
    Uses fuzzy color matching with geometry-aware filtering: tries a
    geometry-compatible sub-legend first (higher tolerance safe because
    cross-type mismatches are impossible), then falls back to full legend.
    """
    compat = {
        "line": ("beam", "small_beam", "wall"),
        "rectangle": ("column",),
    }
    compatible_types = compat.get(geom_type, ())

    # Build geometry-filtered legend (only compatible entry types)
    filtered = {}
    for key, entries in legend.items():
        matches = [e for e in entries if e.element_type in compatible_types]
        if matches:
            filtered[key] = matches

    # Try geometry-compatible match first
    match = _fuzzy_color_match(color_hex, filtered, tolerance)
    if match:
        return match[0]

    # Fallback: try all entries with geometry preference
    match = _fuzzy_color_match(color_hex, legend, tolerance)
    if match:
        for e in match:
            if e.element_type in compatible_types:
                return e
        return match[0]
    return None


def extract_and_classify_shapes(slide, slide_num, legend, scale, floors,
                                phase, slide_w, slide_h, legend_boundary,
                                legend_side, warnings, theme_map=None,
                                tolerance=15, line_tolerance=None):
    """Extract FREEFORM shapes and classify them into structural elements.

    tolerance: fuzzy color matching tolerance (Chebyshev distance) for rectangles.
               Use 150 for table-sourced legends (cell bg ≠ shape fill).
    line_tolerance: override tolerance for LINE shapes. None = use tolerance.
                    Set to 0 for exact match (LINE colors match legend exactly).
    Returns list of element dicts.
    """
    _lt = line_tolerance if line_tolerance is not None else tolerance
    elements = []
    slide_area = slide_w * slide_h

    # Collect all freeforms for column dedup
    filled_rects = []
    outline_rects = []
    lines = []

    for shape in _iter_slide_shapes(slide, (MSO_SHAPE_TYPE.FREEFORM,)):
        w, h = shape.width, shape.height
        left, top = shape.left, shape.top

        # Skip shapes in legend region
        cx = left + w // 2
        if legend_side == "left" and cx < legend_boundary:
            continue
        if legend_side == "right" and cx > legend_boundary:
            continue

        # Skip extremely large shapes (background/border)
        if w > 0 and h > 0 and (w * h) > slide_area * SLIDE_AREA_RATIO_MAX:
            continue

        line_color = _get_shape_color(shape, "line", theme_map)
        fill_color = _get_shape_color(shape, "fill", theme_map)

        if w == 0 or h == 0:
            # Line shape (beam candidate)
            primary_color = line_color or fill_color
            alt_color_line = fill_color if primary_color == line_color else line_color
            lines.append({
                "left": left, "top": top, "width": w, "height": h,
                "color": primary_color,
                "alt_color": alt_color_line,
                "orientation": "H" if h == 0 else "V",
            })
        elif w > 0 and h > 0 and w < 500000 and h < 500000:
            # Small rectangle (column candidate)
            # Only accept 4-vertex freeforms (true rectangles)
            vc = _get_freeform_vertex_count(shape)
            if vc is not None and vc not in COLUMN_REQUIRED_VERTICES:
                continue
            if fill_color:
                filled_rects.append({
                    "left": left, "top": top, "width": w, "height": h,
                    "color": fill_color,
                    "alt_color": line_color,
                })
            elif line_color:
                outline_rects.append({
                    "left": left, "top": top, "width": w, "height": h,
                    "color": line_color,
                    "alt_color": fill_color,
                })

    # ── Deduplicate columns (fill + outline pairs) ──────────────────
    columns = _deduplicate_columns(filled_rects, outline_rects)

    # ── Classify columns ──────────────────────────────────────────────
    if phase in ("phase1", "all"):
        for col in columns:
            entry = _resolve_pptx_legend(col["color"], "rectangle", legend, tolerance)
            # Fallback: try alternate color if primary didn't match column
            if (not entry or entry.element_type != "column") and col.get("alt_color"):
                alt_entry = _resolve_pptx_legend(col["alt_color"], "rectangle", legend, tolerance)
                if alt_entry and alt_entry.element_type == "column":
                    entry = alt_entry
            if not entry or entry.element_type != "column":
                continue

            # Skip shapes too small to be columns
            if col["width"] < MIN_COLUMN_DIM_EMU or col["height"] < MIN_COLUMN_DIM_EMU:
                continue

            cx_emu = col["left"] + col["width"] // 2
            cy_emu = col["top"] + col["height"] // 2
            cx_m, cy_m = _emu_to_m(cx_emu, cy_emu, scale, slide_h)

            # Use legend section if specific; otherwise leave empty
            # (PPT column shapes are not drawn to scale, so dimension-based
            #  section names are unreliable)
            section = entry.section or ""

            elements.append({
                "element_type": "column",
                "x1": cx_m, "y1": cy_m, "x2": cx_m, "y2": cy_m,
                "section": section,
                "floors": list(floors),
                "direction": "",
                "page_num": slide_num,
            })

    # ── Classify lines (beams) ────────────────────────────────────────
    for line in lines:
        color = line["color"]
        if not color:
            continue

        entry = _resolve_pptx_legend(color, "line", legend, _lt)
        # Fallback: try alternate color if primary didn't match
        if not entry and line.get("alt_color"):
            entry = _resolve_pptx_legend(line["alt_color"], "line", legend, _lt)
        if not entry:
            continue

        # Phase filter
        if phase == "phase1" and entry.element_type == "small_beam":
            continue
        if phase == "phase2" and entry.element_type != "small_beam":
            continue

        if line["orientation"] == "H":
            # Horizontal line: Y-direction beam
            x1_emu = line["left"]
            x2_emu = line["left"] + line["width"]
            y_emu = line["top"]
            x1_m, y1_m = _emu_to_m(x1_emu, y_emu, scale, slide_h)
            x2_m, y2_m = _emu_to_m(x2_emu, y_emu, scale, slide_h)
        else:
            # Vertical line: X-direction beam
            x_emu = line["left"]
            y1_emu = line["top"]
            y2_emu = line["top"] + line["height"]
            x1_m, y1_m = _emu_to_m(x_emu, y1_emu, scale, slide_h)
            x2_m, y2_m = _emu_to_m(x_emu, y2_emu, scale, slide_h)

        direction = _direction_of(x1_m, y1_m, x2_m, y2_m)

        elem = {
            "element_type": entry.element_type,
            "x1": x1_m, "y1": y1_m, "x2": x2_m, "y2": y2_m,
            "section": entry.section or "",
            "floors": list(floors),
            "direction": direction,
            "page_num": slide_num,
        }
        if entry.element_type == "wall" and entry.is_diaphragm:
            elem["is_diaphragm_wall"] = True
        if not entry.section:
            elem["section_uncertain"] = True
        elements.append(elem)

    return elements


def _deduplicate_columns(filled_rects, outline_rects):
    """Merge fill + outline FREEFORM pairs into single column entries.

    Two freeforms at the same position (within tolerance) represent one column:
    one with solid fill, one with outline only.
    """
    # Start with filled rects as the primary set
    columns = list(filled_rects)

    # Mark outline rects that overlap with a filled rect
    used_fills = set()
    for orect in outline_rects:
        ox, oy = orect["left"], orect["top"]
        matched = False
        for i, frect in enumerate(filled_rects):
            if i in used_fills:
                continue
            fx, fy = frect["left"], frect["top"]
            if (abs(ox - fx) < COLUMN_DEDUP_TOLERANCE_EMU and
                    abs(oy - fy) < COLUMN_DEDUP_TOLERANCE_EMU):
                matched = True
                used_fills.add(i)
                break
        if not matched:
            # Outline-only rect without matching fill — still a column
            columns.append(orect)

    return columns


# ─── Summary Report ──────────────────────────────────────────────────────────

def print_summary(output: dict):
    """Print human-readable summary to stdout."""
    meta = output["_metadata"]
    print(f"\n{'='*60}")
    print(f"pptx_to_elements.py — Summary")
    print(f"{'='*60}")
    print(f"Input:  {meta['input_file']}")
    print(f"Phase:  {meta['phase']}")

    if meta.get("scale_details"):
        for sd in meta["scale_details"]:
            if isinstance(sd, dict):
                print(f"Scale:  Slide {sd.get('slide', '?')}: "
                      f"{sd.get('emu_per_meter', '?')} EMU/m "
                      f"({sd.get('num_measurements', 0)} measurements)")

    print(f"\nPer-slide element counts:")
    for pn, stats in meta.get("per_page_stats", {}).items():
        floors = meta.get("page_floors", {}).get(pn, [])
        floor_desc = (f"{floors[0]}~{floors[-1]}" if len(floors) > 1
                      else (floors[0] if floors else "?"))
        parts = [f"{k}={v}" for k, v in stats.items() if v > 0]
        print(f"  Slide {pn} ({floor_desc}): {', '.join(parts) or 'none'}")

    totals = meta.get("totals", {})
    print(f"\nTotals (after dedup):")
    for k, v in totals.items():
        print(f"  {k}: {v}")

    secs = output.get("sections", {})
    if secs.get("frame"):
        print(f"\nFrame sections: {', '.join(secs['frame'])}")
    if secs.get("wall"):
        print(f"Wall thicknesses (cm): {', '.join(str(w) for w in secs['wall'])}")

    if meta.get("warnings"):
        print(f"\nWARNINGS ({len(meta['warnings'])}):")
        for w in meta["warnings"]:
            print(f"  ! {w}")

    print(f"{'='*60}\n")


# ─── Slide Listing ───────────────────────────────────────────────────────────

def list_slides(prs):
    """Print slide overview with shape type counts."""
    print(f"\nSlides: {len(prs.slides)}")
    print(f"Slide size: {prs.slide_width} x {prs.slide_height} EMU "
          f"({prs.slide_width / 914400:.1f} x {prs.slide_height / 914400:.1f} in)")
    print()
    for i, slide in enumerate(prs.slides):
        counts = {}
        for shape in _iter_slide_shapes(slide):
            stype = str(shape.shape_type).split(".")[-1].split("(")[0].strip()
            counts[stype] = counts.get(stype, 0) + 1
        parts = [f"{k}={v}" for k, v in sorted(counts.items())]
        # Find floor label texts (broadened to match ranges and more shape types)
        floor_texts = []
        for shape in _iter_text_shapes(slide):
            t = shape.text_frame.text.strip()
            for line in t.split("\n"):
                for m in FLOOR_LABEL_RE.finditer(line.strip()):
                    val = m.group(1).strip()
                    if val not in floor_texts:
                        floor_texts.append(val)
        floor_info = f" (floors: {', '.join(floor_texts)})" if floor_texts else ""
        print(f"  Slide {i+1}: {', '.join(parts)}{floor_info}")
    print()


# ─── Main Pipeline ───────────────────────────────────────────────────────────

def _parse_color_map(color_map_args):
    """Parse --color-map arguments into legend entries.

    Args:
        color_map_args: list of strings like "FF0000=column:C100X100"

    Returns:
        dict[color_hex, list[LegendEntry]]
    """
    if not color_map_args:
        return {}
    result = {}
    _PREFIX_MAP = {
        "column": "C", "beam": "B", "wall": "W",
        "small_beam": "SB",
    }
    for entry_str in color_map_args:
        entry_str = entry_str.strip()
        if "=" not in entry_str:
            print(f"  WARNING: invalid --color-map entry (no '='): {entry_str}")
            continue
        color_part, mapping_part = entry_str.split("=", 1)
        color_hex = color_part.strip().upper()
        if len(color_hex) != 6:
            print(f"  WARNING: invalid color hex '{color_hex}' in --color-map")
            continue
        # Parse type:section
        if ":" in mapping_part:
            elem_type, section = mapping_part.split(":", 1)
        else:
            elem_type = mapping_part
            section = None
        elem_type = elem_type.strip().lower()
        if section:
            section = section.strip()
        if elem_type not in ("column", "beam", "wall", "small_beam"):
            print(f"  WARNING: unknown element type '{elem_type}' in --color-map")
            continue
        # Derive prefix from section or element type
        prefix = ""
        if section:
            m = re.match(r"^([A-Z]+)", section)
            if m:
                prefix = m.group(1)
        if not prefix:
            prefix = _PREFIX_MAP.get(elem_type, "")
        is_diaphragm = section and "連續壁" in section
        rgb = [int(color_hex[i:i+2], 16) for i in (0, 2, 4)]
        le = LegendEntry(
            element_type=elem_type,
            section=section if section else None,
            color_name=f"#{color_hex}",
            color_rgb=rgb,
            specificity=1 if section else 0,
            is_diaphragm=is_diaphragm,
            prefix=prefix,
            label=f"[manual] {elem_type}:{section or 'generic'}",
        )
        result.setdefault(color_hex, []).append(le)
    return result


def process(prs, page_floors, phase, crop=False, crop_dir=None,
            manual_scale=None, color_map=None, slides_info_dir=None):
    """Main processing pipeline: PPTX → elements dict.

    Args:
        prs: python-pptx Presentation object.
        page_floors: {slide_num: [floor_names]}.
        phase: "phase1", "phase2", or "all".
        crop: Whether to extract PNG images.
        crop_dir: Directory for extracted PNGs.
        color_map: dict[color_hex, list[LegendEntry]] from --color-map.

    Returns:
        Complete output dict ready for JSON serialization.
    """
    warnings = []
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Build theme color map once for the whole presentation
    theme_map = _build_theme_color_map(prs)
    if theme_map:
        print(f"  Theme colors resolved: {len(theme_map)} entries")

    # Validate slide references
    num_slides = len(prs.slides)
    for sn in page_floors:
        if sn < 1 or sn > num_slides:
            warnings.append(
                f"Slide {sn} referenced in --page-floors but only {num_slides} slides exist"
            )

    all_elements = []
    per_page_stats = {}
    scale_details = []
    legend_validations = []
    legend_diagnostics_all = []

    # Pre-compute scales for all requested slides; allow fallback
    slide_scales = {}
    if manual_scale:
        # Manual scale overrides auto-detection for all slides
        for sn in sorted(page_floors.keys()):
            if sn < 1 or sn > num_slides:
                continue
            slide_scales[sn] = (manual_scale, {
                "slide": sn, "emu_per_meter": round(manual_scale, 1),
                "num_measurements": 0, "manual": True
            })
        print(f"  Using manual scale: {manual_scale:.1f} EMU/m for all slides")
    else:
        for sn in sorted(page_floors.keys()):
            if sn < 1 or sn > num_slides:
                continue
            scale_val, scale_info = detect_slide_scale(prs.slides[sn - 1], sn)
            if scale_val is not None:
                slide_scales[sn] = (scale_val, scale_info)
            else:
                # Fallback: estimate from shape extent
                est_scale, est_info = _estimate_scale_from_shapes(prs.slides[sn - 1], sn)
                if est_scale is not None:
                    slide_scales[sn] = (est_scale, est_info)
                    print(f"  Slide {sn}: using estimated scale {est_scale:.0f} EMU/m (no measurement texts)")

    # Fallback scale: use median of all detected scales (or manual scale)
    fallback_scale = None
    if slide_scales:
        all_scales = sorted(v[0] for v in slide_scales.values())
        fallback_scale = all_scales[len(all_scales) // 2]

    for sn, floors in sorted(page_floors.items()):
        if sn < 1 or sn > num_slides:
            continue
        slide = prs.slides[sn - 1]  # 1-based to 0-based
        print(f"\n--- Processing Slide {sn} ({', '.join(floors)}) ---")

        # 1) Compute floor label (needed for crop + JSON)
        floor_label = (f"{floors[0]}~{floors[-1]}" if len(floors) > 1
                       else floors[0])

        # 2) Extract PNG FIRST (always, regardless of scale)
        if crop:
            actual_crop_dir = (str(Path(slides_info_dir) / floor_label / "screenshots")) if slides_info_dir else crop_dir
            if actual_crop_dir:
                extract_base_png(slide, sn, floor_label, actual_crop_dir)

        # 3) Detect scale (per-slide, with fallback)
        if sn in slide_scales:
            scale, scale_info = slide_scales[sn]
            if isinstance(scale_info, dict) and scale_info.get("estimated"):
                print(f"  Scale: {scale:.0f} EMU/m (ESTIMATED from shape extent \u2014 affine_calibrate will correct)")
            else:
                print(f"  Scale: {scale:.1f} EMU/m ({scale_info.get('num_measurements', 0) if isinstance(scale_info, dict) else 0} measurements)")
        elif fallback_scale:
            scale = fallback_scale
            scale_info = {"slide": sn, "emu_per_meter": round(scale, 1),
                          "num_measurements": 0, "fallback": True}
            warnings.append(
                f"Slide {sn}: no measurement texts, using fallback scale "
                f"{scale:.1f} EMU/m from other slides"
            )
            print(f"  Scale: {scale:.1f} EMU/m (FALLBACK from other slides)")
        else:
            warnings.append(f"Slide {sn}: no scale available, skipping")
            print(f"  ERROR: no scale available")
            continue
        scale_details.append(scale_info)

        # 4) Parse legend (table-first, shape fallback)
        legend, legend_diag, legend_boundary, legend_side = parse_legend(
            slide, slide_w, slide_h, phase=phase, theme_map=theme_map)
        legend_source = legend_diag.get("legend_source", "shape")
        # Inject manual color-map entries into legend
        if color_map:
            for cm_color, cm_entries in color_map.items():
                if cm_color not in legend:
                    legend[cm_color] = list(cm_entries)
                else:
                    # Append manual entries, avoiding duplicates
                    existing_types = {(e.element_type, e.section) for e in legend[cm_color]}
                    for cme in cm_entries:
                        if (cme.element_type, cme.section) not in existing_types:
                            legend[cm_color].append(cme)
        if not legend:
            warnings.append(f"Slide {sn}: no legend detected, skipping")
            print(f"  WARNING: no legend detected")
            continue
        print(f"  Legend ({legend_source}): {len(legend)} colors → "
              f"{sum(len(v) for v in legend.values())} entries")
        for color, entries in legend.items():
            for e in entries:
                print(f"    #{color} → {e.element_type}: {e.section or '(generic)'} [{e.label}]")
        if legend_diag.get("pass2_labels"):
            print(f"  Legend Pass 2 (relaxed): {len(legend_diag['pass2_labels'])} additional labels matched")
        if legend_diag.get("unmatched_labels"):
            for ul in legend_diag["unmatched_labels"]:
                warnings.append(f"Slide {sn}: label '{ul['label']}' ({ul['type']}) unmatched to any swatch")
        if legend_diag.get("unmatched_swatches"):
            for us in legend_diag["unmatched_swatches"]:
                warnings.append(f"Slide {sn}: swatch #{us['color']} has no legend label")

        # Table legend color remap diagnostics
        if legend_source == "table":
            # Collect all unique shape colors in the drawing area
            shape_fill_colors = set()
            for shape in _iter_slide_shapes(slide, (MSO_SHAPE_TYPE.FREEFORM,)):
                cx = shape.left + shape.width // 2
                if legend_side == "left" and cx < legend_boundary:
                    continue
                if legend_side == "right" and cx > legend_boundary:
                    continue
                lc = _get_shape_color(shape, "line", theme_map)
                fc = _get_shape_color(shape, "fill", theme_map)
                if lc and lc != "FFFFFF":
                    shape_fill_colors.add(lc)
                if fc and fc != "FFFFFF":
                    shape_fill_colors.add(fc)
            if shape_fill_colors:
                print(f"  Table legend color mapping check (tolerance=150):")
                for lcolor in legend:
                    lr = int(lcolor[0:2], 16)
                    lg_ = int(lcolor[2:4], 16)
                    lb = int(lcolor[4:6], 16)
                    best_sc, best_d = None, float('inf')
                    for sc in shape_fill_colors:
                        sr = int(sc[0:2], 16)
                        sg = int(sc[2:4], 16)
                        sb = int(sc[4:6], 16)
                        d = max(abs(lr-sr), abs(lg_-sg), abs(lb-sb))
                        if d < best_d:
                            best_d = d
                            best_sc = sc
                    label = legend[lcolor][0].section or legend[lcolor][0].label
                    if best_d == 0:
                        print(f"    #{lcolor} ({label}) → closest shape: #{best_sc} (exact)")
                    else:
                        print(f"    #{lcolor} ({label}) → closest shape: #{best_sc} (dist={best_d})")

        # 4) Legend validation report
        legend_tolerance = 150 if legend_source == "table" else 15
        lv = _validate_legend(slide, sn, legend, slide_w, slide_h,
                              legend_side, legend_boundary, warnings,
                              theme_map=theme_map, tolerance=legend_tolerance)
        _print_legend_validation(lv)
        legend_validations.append(lv)
        legend_diagnostics_all.append({"slide": sn, **legend_diag})

        # 5) Extract & classify shapes
        rect_tolerance = 150 if legend_source == "table" else 15
        slide_elems = extract_and_classify_shapes(
            slide, sn, legend, scale, floors, phase,
            slide_w, slide_h, legend_boundary, legend_side, warnings,
            theme_map=theme_map, tolerance=rect_tolerance, line_tolerance=0)
        all_elements.extend(slide_elems)

        # Per-slide stats
        stats = {"beams": 0, "columns": 0, "walls": 0, "small_beams": 0}
        for e in slide_elems:
            key = e["element_type"] + "s"
            if key in stats:
                stats[key] += 1
        per_page_stats[str(sn)] = stats
        print(f"  Elements: {', '.join(f'{k}={v}' for k, v in stats.items() if v > 0) or 'none'}")

        # Per-slide JSON output (when --slides-info-dir is specified)
        if slides_info_dir:
            slide_output = {
                "_metadata": {
                    "slide_num": sn,
                    "floors": list(floors),
                    "floor_label": floor_label,
                    "scale": scale_info,
                    "legend": {
                        color: [{"element_type": e.element_type,
                                 "section": e.section,
                                 "label": e.label}
                                for e in entries]
                        for color, entries in legend.items()
                    },
                    "legend_validation": lv,
                    "stats": stats,
                },
                "columns": [e for e in slide_elems if e["element_type"] == "column"],
                "beams": [e for e in slide_elems if e["element_type"] == "beam"],
                "walls": [e for e in slide_elems if e["element_type"] == "wall"],
                "small_beams": [e for e in slide_elems if e["element_type"] == "small_beam"],
            }
            # Phase 2 per-slide files use sb_ prefix to avoid collision with Phase 1
            if phase == "phase2":
                slide_json_path = Path(slides_info_dir) / floor_label / f"sb_{floor_label}.json"
            else:
                slide_json_path = Path(slides_info_dir) / floor_label / f"{floor_label}.json"
            slide_json_path.parent.mkdir(parents=True, exist_ok=True)
            with open(slide_json_path, "w", encoding="utf-8") as f:
                json.dump(slide_output, f, ensure_ascii=False, indent=2)
            print(f"  Saved per-slide JSON: {slide_json_path}")

    # When --slides-info-dir is set, skip merge and return summary
    if slides_info_dir:
        slide_files = sorted(Path(slides_info_dir).glob("*/*.json"))
        print(f"\n--- Per-slide JSON output complete ---")
        print(f"  Directory: {slides_info_dir}")
        print(f"  Files: {len(slide_files)}")
        for sf in slide_files:
            print(f"    {sf.name}")
        return {
            "_metadata": {
                "source": "pptx_to_elements.py",
                "mode": "per-slide",
                "slides_info_dir": str(slides_info_dir),
                "input_file": str(getattr(prs, '_pptx_path', '')),
                "generated_at": datetime.now(timezone.utc)
                                .strftime("%Y-%m-%dT%H:%M:%S"),
                "phase": phase,
                "page_floors": {str(k): v for k, v in page_floors.items()},
                "per_page_stats": per_page_stats,
                "warnings": warnings,
            },
            "slide_files": [str(sf) for sf in slide_files],
        }

    # Snap walls to nearest beam axis (post-processing)
    _snap_walls_to_beams(all_elements, tolerance=1.0, warnings=warnings)

    # Group & dedup
    grouped = group_and_dedup(all_elements)
    sections = collect_sections(grouped)

    # Section coverage diagnostics
    total_elems = (len(grouped["columns"]) + len(grouped["beams"]) +
                   len(grouped["walls"]) + len(grouped["small_beams"]))
    empty_section_count = sum(
        1 for cat in ("columns", "beams", "walls", "small_beams")
        for elem in grouped[cat]
        if not elem.get("section")
    )
    if empty_section_count > 0 and total_elems > 0:
        pct = empty_section_count / total_elems * 100
        msg = (f"Section coverage: {empty_section_count}/{total_elems} "
               f"elements ({pct:.0f}%) have empty section")
        if pct > 20:
            warnings.append(
                f"WARNING: 大量元素無斷面 — {msg}. "
                f"Check Legend color matching (fill vs outline, tolerance).")
        else:
            warnings.append(msg)

    # Build output
    output = {
        "_metadata": {
            "source": "pptx_to_elements.py",
            "input_file": str(getattr(prs, '_pptx_path', '')),
            "generated_at": datetime.now(timezone.utc)
                            .strftime("%Y-%m-%dT%H:%M:%S"),
            "phase": phase,
            "page_floors": {str(k): v for k, v in page_floors.items()},
            "scale_details": scale_details,
            "per_page_stats": per_page_stats,
            "totals": {
                "columns": len(grouped["columns"]),
                "beams": len(grouped["beams"]),
                "walls": len(grouped["walls"]),
                "small_beams": len(grouped["small_beams"]),
            },
            "legend_validation": legend_validations,
            "legend_diagnostics": legend_diagnostics_all,
            "warnings": warnings,
        },
        "columns": grouped["columns"],
        "beams": grouped["beams"],
        "walls": grouped["walls"],
        "small_beams": grouped["small_beams"],
        "sections": sections,
    }
    return output


# ─── CLI ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="PPT structural plan → elements JSON"
    )
    parser.add_argument(
        "--input", "-i", required=True,
        help="Path to .pptx file")
    parser.add_argument(
        "--output", "-o",
        help="Output JSON file path (elements.json or sb_elements.json)")
    parser.add_argument(
        "--page-floors",
        help='Slide-to-floor mapping, e.g. "1=B3F, 3=1F~2F, 4=3F~14F"')
    parser.add_argument(
        "--phase", choices=["phase1", "phase2", "all"], default="all",
        help="phase1=beams+columns+walls, phase2=small_beams, all=both")
    parser.add_argument(
        "--crop", action="store_true",
        help="Extract PNG base images from slides")
    parser.add_argument(
        "--crop-dir",
        help="Directory for extracted PNGs")
    parser.add_argument(
        "--dry-run", action="store_true",
        help="Show summary without writing output file")
    parser.add_argument(
        "--list-slides", action="store_true",
        help="List slides and exit")
    parser.add_argument(
        "--scan-floors", action="store_true",
        help="Scan all slides for floor labels and print suggested --page-floors, then exit")
    parser.add_argument(
        "--auto-floors", action="store_true",
        help="Auto-detect floor labels and use them (no --page-floors needed)")
    parser.add_argument(
        "--confirm-floors", action="store_true",
        help="Scan floors, display detections with confidence, and prompt for confirmation")
    parser.add_argument(
        "--scale", type=float, default=None,
        help="Manual scale in EMU/m (overrides auto-detection from measurement texts)")
    parser.add_argument(
        "--color-map", nargs="*", default=None,
        help='Manual color-to-element mappings for unmatched colors. '
             'Format: "RRGGBB=type:SECTION" e.g. "FF0000=column:C100X100" '
             '"FF0000=wall:W90" "FFFF00=beam:B60X80". '
             'Multiple entries per color allowed.')
    parser.add_argument(
        "--color-tolerance", type=int, default=None,
        help="Override fuzzy color matching tolerance (default: 15 per RGB channel)")
    parser.add_argument(
        "--slides-info-dir",
        help="Directory for per-slide JSON output (e.g., '結構配置圖/SLIDES INFO'). "
             "When set, writes per-slide JSONs instead of merged output.")

    args = parser.parse_args()

    # Validate mutually exclusive options
    floor_opts = sum([bool(args.auto_floors), bool(args.page_floors), bool(args.confirm_floors)])
    if floor_opts > 1:
        print("ERROR: --auto-floors, --confirm-floors, and --page-floors are mutually exclusive")
        sys.exit(1)

    # Apply color tolerance override if specified
    global _COLOR_TOLERANCE_OVERRIDE
    if args.color_tolerance is not None:
        _COLOR_TOLERANCE_OVERRIDE = args.color_tolerance

    # Load PPTX
    input_path = Path(args.input)
    if not input_path.exists():
        print(f"ERROR: File not found: {input_path}")
        sys.exit(1)

    prs = Presentation(str(input_path))
    prs._pptx_path = str(input_path)
    print(f"Loaded: {input_path.name} ({len(prs.slides)} slides)")

    # List mode
    if args.list_slides:
        list_slides(prs)
        return

    # Scan-floors mode: detect and print with confidence, then exit
    if args.scan_floors:
        detected = scan_floors(prs)
        if not detected:
            print("No floor labels detected on any slide.")
        else:
            print(format_scan_floors_output(detected))
        return

    # Confirm-floors mode: scan, show confidence, prompt user
    if args.confirm_floors:
        detected = scan_floors(prs)
        if not detected:
            print("No floor labels detected on any slide.")
            sys.exit(1)
        print(format_scan_floors_output(detected))
        try:
            answer = input("Confirm? [Y/n]: ").strip().lower()
        except (EOFError, KeyboardInterrupt):
            answer = "n"
        if answer and answer != "y":
            print("Aborted by user.")
            sys.exit(1)
        # Use confirmed labels as page_floors
        page_floors = {}
        for sn, entries in detected.items():
            labels = [e["label"] for e in entries if e["confidence"] != "low"]
            if not labels:
                labels = [e["label"] for e in entries]
            all_floors = []
            for label in labels:
                all_floors.extend(expand_floor_range(label))
            seen = set()
            deduped = []
            for f in all_floors:
                if f not in seen:
                    seen.add(f)
                    deduped.append(f)
            if deduped:
                page_floors[sn] = deduped
        print(f"Confirmed floor mappings: {len(page_floors)} slides")

    # Auto-floors mode: detect and use
    elif args.auto_floors:
        detected = scan_floors_labels(prs)
        if not detected:
            print("ERROR: --auto-floors found no floor labels on any slide")
            sys.exit(1)
        # Convert detected labels to page_floors format
        page_floors = {}
        for sn, labels in detected.items():
            all_floors = []
            for label in labels:
                all_floors.extend(expand_floor_range(label))
            # Deduplicate
            seen = set()
            deduped = []
            for f in all_floors:
                if f not in seen:
                    seen.add(f)
                    deduped.append(f)
            page_floors[sn] = deduped
        print(f"Auto-detected floor mappings: {len(page_floors)} slides")
        for sn in sorted(page_floors.keys()):
            floors = page_floors[sn]
            print(f"  Slide {sn}: {', '.join(floors)}")
    elif args.page_floors:
        # Parse page-floors
        page_floors = parse_page_floors(args.page_floors)
        if not page_floors:
            print("ERROR: No valid slide-floor mappings in --page-floors")
            sys.exit(1)
        print(f"Slide-floor mappings: {len(page_floors)} slides")
    else:
        print("ERROR: --page-floors, --auto-floors, or --confirm-floors is required for processing")
        sys.exit(1)

    # Require --output for processing (unless --slides-info-dir is used)
    if not args.output and not args.slides_info_dir:
        print("ERROR: --output or --slides-info-dir is required for processing")
        sys.exit(1)

    # Parse color-map if provided
    manual_color_map = _parse_color_map(args.color_map) if args.color_map else None
    if manual_color_map:
        print(f"Manual color-map: {sum(len(v) for v in manual_color_map.values())} entries")
        for c, entries in manual_color_map.items():
            for e in entries:
                print(f"  #{c} → {e.element_type}: {e.section or '(generic)'}")

    # Process
    output = process(prs, page_floors, args.phase,
                     crop=args.crop, crop_dir=args.crop_dir,
                     manual_scale=args.scale, color_map=manual_color_map,
                     slides_info_dir=args.slides_info_dir)

    # Summary and output
    if args.slides_info_dir:
        # Per-slide mode: files already written in process(), skip merged output
        print(f"\nPer-slide output complete. Files in: {args.slides_info_dir}")
    else:
        print_summary(output)
        if args.dry_run:
            print("[DRY RUN] No output written.")
        else:
            out_path = Path(args.output)
            out_path.parent.mkdir(parents=True, exist_ok=True)
            with open(out_path, "w", encoding="utf-8") as f:
                json.dump(output, f, ensure_ascii=False, indent=2)
            print(f"Output written to: {out_path}")


if __name__ == "__main__":
    main()
