"""
Generate PPT: 用 AI Agent 打造智能工作流
Dark theme with bright accents, 16 slides
Concept-first (70%) + Application (30%)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)       # Dark navy background
BG_CARD = RGBColor(0x24, 0x24, 0x3E)       # Card background (slightly lighter)
BG_TERMINAL = RGBColor(0x0D, 0x0D, 0x1A)   # Terminal background
ACCENT_ORANGE = RGBColor(0xFF, 0x6B, 0x35)  # Orange accent
ACCENT_CYAN = RGBColor(0x00, 0xD2, 0xFF)    # Cyan accent
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x76)   # Green for positive
ACCENT_RED = RGBColor(0xFF, 0x4D, 0x4D)     # Red for negative
ACCENT_PURPLE = RGBColor(0xB4, 0x7E, 0xFF)  # Purple accent (new)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x88, 0x88, 0x99)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)


def set_slide_bg(slide, color=BG_DARK):
    """Set slide background to solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft JhengHei"):
    """Add a textbox with styled text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=18,
                  color=WHITE, bold=False, spacing=1.2, alignment=PP_ALIGN.LEFT,
                  font_name="Microsoft JhengHei"):
    """Add textbox with multiple lines (list of (text, color, bold, font_size) or str)."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(line, str):
            p.text = line
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.bold = bold
        elif isinstance(line, tuple):
            txt = line[0]
            clr = line[1] if len(line) > 1 else color
            bld = line[2] if len(line) > 2 else bold
            fsz = line[3] if len(line) > 3 else font_size
            p.text = txt
            p.font.size = Pt(fsz)
            p.font.color.rgb = clr
            p.font.bold = bld

        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * (spacing - 1) + 4)

    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD,
                     border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_arrow(slide, left, top, width, height, color=ACCENT_ORANGE):
    """Add a right arrow shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_down_arrow(slide, left, top, width, height, color=ACCENT_ORANGE):
    """Add a down arrow shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_title_bar(slide, title_text):
    """Add a consistent title bar at the top of the slide."""
    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(0.45), Inches(1.2), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_ORANGE
    line.line.fill.background()

    add_textbox(slide, 0.6, 0.55, 8.5, 0.8, title_text,
                font_size=32, color=WHITE, bold=True)


def add_section_label(slide, text, y=0.45):
    """Add section label (概念篇/應用篇) at top-right."""
    add_textbox(slide, 10.5, y, 2.5, 0.4,
                text, font_size=14, color=MID_GRAY,
                alignment=PP_ALIGN.RIGHT)


def add_layer_badge(slide, layer_num, x=12.0, y=0.45):
    """Add a layer number badge at top-right."""
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x), Inches(y), Inches(0.5), Inches(0.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_ORANGE
    circle.line.fill.background()
    add_textbox(slide, x, y + 0.05, 0.5, 0.4,
                str(layer_num), font_size=20, color=BG_DARK, bold=True,
                alignment=PP_ALIGN.CENTER)


def add_page_number(slide, num, total=16):
    """Add page number at bottom-right."""
    add_textbox(slide, 11.5, 7.1, 1.5, 0.3,
                f"{num} / {total}", font_size=12, color=MID_GRAY,
                alignment=PP_ALIGN.RIGHT)


def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ══════════════════════════════════════════════
    # SLIDE 1: Cover
    # ══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # Top bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_ORANGE
    bar.line.fill.background()

    # Title
    add_textbox(slide, 1.5, 2.0, 10, 1.2,
                "用 AI Agent 打造智能工作流",
                font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Subtitle
    add_textbox(slide, 1.5, 3.3, 10, 0.8,
                "從概念到落地——分層架構的力量",
                font_size=28, color=ACCENT_CYAN, bold=False, alignment=PP_ALIGN.CENTER)

    # Divider
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.0), Inches(4.3), Inches(3.333), Inches(0.04)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = MID_GRAY
    divider.line.fill.background()

    # Date / info
    add_textbox(slide, 1.5, 4.6, 10, 0.6,
                "2026 | Structural AI Solutions",
                font_size=20, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # Bottom bar
    bar2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.42), Inches(13.333), Inches(0.08)
    )
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = ACCENT_ORANGE
    bar2.line.fill.background()

    # ══════════════════════════════════════════════
    # SLIDE 2: What is AI Agent?
    # ══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "從 ChatBot 到 Agent — AI 的進化")
    add_section_label(slide, "概念篇")
    add_page_number(slide, 2)

    # ChatBot card (left)
    add_rounded_rect(slide, 0.8, 1.6, 5.5, 5.2, BG_CARD, ACCENT_RED)
    add_textbox(slide, 0.9, 1.8, 5.3, 0.6,
                "ChatBot", font_size=36, color=ACCENT_RED, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0.9, 2.5, 5.3, 0.4,
                "被動問答", font_size=20, color=MID_GRAY,
                alignment=PP_ALIGN.CENTER)

    chatbot_items = [
        ("互動模式", "你問一句、它答一句"),
        ("工具能力", "只能產生文字"),
        ("協作方式", "單一對話"),
        ("結果", "一段回答"),
    ]
    for i, (label, desc) in enumerate(chatbot_items):
        y = 3.2 + i * 0.75
        add_textbox(slide, 1.1, y, 2.0, 0.4,
                    label, font_size=16, color=ACCENT_RED, bold=True)
        add_textbox(slide, 3.1, y, 3.0, 0.4,
                    desc, font_size=16, color=LIGHT_GRAY)

    # Arrow in the middle
    add_arrow(slide, 6.4, 3.5, 0.6, 0.5, ACCENT_ORANGE)
    add_textbox(slide, 6.2, 2.8, 1.0, 0.5,
                "進化", font_size=16, color=ACCENT_ORANGE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Agent card (right)
    add_rounded_rect(slide, 7.1, 1.6, 5.5, 5.2, BG_CARD, ACCENT_GREEN)
    add_textbox(slide, 7.2, 1.8, 5.3, 0.6,
                "AI Agent", font_size=36, color=ACCENT_GREEN, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 7.2, 2.5, 5.3, 0.4,
                "主動執行", font_size=20, color=MID_GRAY,
                alignment=PP_ALIGN.CENTER)

    agent_items = [
        ("互動模式", "給任務，自己拆解步驟完成"),
        ("工具能力", "讀寫檔案、執行程式、操作工具"),
        ("協作方式", "多 Agent 分工合作"),
        ("結果", "實際成果 — 檔案、模型、報告"),
    ]
    for i, (label, desc) in enumerate(agent_items):
        y = 3.2 + i * 0.75
        add_textbox(slide, 7.3, y, 2.0, 0.4,
                    label, font_size=16, color=ACCENT_GREEN, bold=True)
        add_textbox(slide, 9.3, y, 3.1, 0.4,
                    desc, font_size=16, color=WHITE)

    # ══════════════════════════════════════════════
    # SLIDE 3: Claude Code
    # ══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Claude Code：讓 Agent 落地的平台")
    add_section_label(slide, "概念篇")
    add_page_number(slide, 3)

    # Left: description
    add_multiline(slide, 0.8, 1.6, 5.5, 2.0, [
        ("命令列工具（CLI）", WHITE, True, 24),
        ("在終端機直接與 AI 對話、下達任務", LIGHT_GRAY, False, 18),
        ("不是通用聊天工具，是專業工作平台", ACCENT_CYAN, True, 18),
    ], font_size=18)

    # Four capabilities - 2x2 grid
    capabilities = [
        ("📂", "讀寫任何檔案", ACCENT_CYAN),
        ("💻", "執行程式碼", ACCENT_GREEN),
        ("🔧", "操作外部工具與軟體", ACCENT_ORANGE),
        ("🧩", "自訂指令 · Agent · 技能", ACCENT_PURPLE),
    ]

    for i, (icon, desc, accent) in enumerate(capabilities):
        col = i % 2
        row = i // 2
        x = 0.8 + col * 3.0
        y = 3.8 + row * 1.5
        add_rounded_rect(slide, x, y, 2.8, 1.2, BG_CARD, accent)
        add_textbox(slide, x + 0.1, y + 0.1, 0.6, 0.5,
                    icon, font_size=28, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.7, y + 0.15, 2.0, 1.0,
                    desc, font_size=17, color=WHITE, bold=True)

    # Right: terminal mockup
    add_rounded_rect(slide, 7.0, 1.6, 5.8, 5.3, BG_TERMINAL, ACCENT_CYAN)
    # Terminal title bar
    add_rounded_rect(slide, 7.0, 1.6, 5.8, 0.45, RGBColor(0x2D, 0x2D, 0x4E))
    add_textbox(slide, 7.2, 1.65, 3.0, 0.35,
                "Claude Code", font_size=14, color=LIGHT_GRAY,
                font_name="Consolas")

    terminal_lines = [
        ("> 請幫我讀取結構平面圖並建立模型", WHITE, False, 16),
        ("", WHITE, False, 8),
        ("正在讀取 A21_3F.pdf ...", ACCENT_CYAN, False, 14),
        ("辨識到 Grid: X 5軸, Y 5軸", ACCENT_GREEN, False, 14),
        ("辨識到 柱: 12 種斷面", ACCENT_GREEN, False, 14),
        ("辨識到 梁: 8 種斷面", ACCENT_GREEN, False, 14),
        ("", WHITE, False, 8),
        ("生成 model_config.json ...", ACCENT_ORANGE, False, 14),
        ("啟動 Golden Scripts 建模 ...", ACCENT_ORANGE, False, 14),
        ("", WHITE, False, 8),
        ("Step 01/11: 材料定義 ✓", ACCENT_GREEN, False, 14),
        ("Step 02/11: 斷面展開 ✓", ACCENT_GREEN, False, 14),
        ("Step 03/11: Grid+樓層 ✓", ACCENT_GREEN, False, 14),
        ("...", MID_GRAY, False, 14),
    ]
    add_multiline(slide, 7.3, 2.2, 5.3, 4.5, terminal_lines,
                  font_size=14, font_name="Consolas", spacing=1.0)

    # ══════════════════════════════════════════════
    # SLIDE 4: Architecture Overview (Pyramid)
    # ══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "六層架構 — 從知識到執行")
    add_section_label(slide, "概念篇")
    add_page_number(slide, 4)

    # Pyramid layers (bottom to top) - using stacked rectangles with decreasing width
    layers = [
        ("第 1 層：專案說明書",     "AI 工作手冊",       ACCENT_CYAN,    11.0, 0.80),
        ("第 2 層：工作流命令",     "一鍵啟動流程",      ACCENT_CYAN,     9.5, 0.80),
        ("第 3 層：AI Agent 團隊", "多專家分工",        ACCENT_ORANGE,   8.0, 0.80),
        ("第 4 層：專業技能模組",    "可複用知識庫",     ACCENT_ORANGE,   6.5, 0.80),
        ("第 5 層：確定性執行腳本",  "同輸入 = 同輸出",  ACCENT_GREEN,    5.0, 0.80),
        ("第 6 層：自動化驗證",     "每次自動品質檢查",   ACCENT_GREEN,   3.5, 0.80),
    ]

    base_y = 6.0  # bottom of pyramid
    center_x = 6.666  # center of slide

    for i, (name, desc, accent, width, height) in enumerate(layers):
        y = base_y - (i + 1) * (height + 0.08)
        x = center_x - width / 2
        add_rounded_rect(slide, x, y, width, height, BG_CARD, accent)
        add_textbox(slide, x + 0.2, y + 0.05, width * 0.55, height - 0.1,
                    name, font_size=18, color=WHITE, bold=True)
        add_textbox(slide, x + width * 0.58, y + 0.05, width * 0.4, height - 0.1,
                    desc, font_size=16, color=accent,
                    alignment=PP_ALIGN.RIGHT)

    # Bottom label
    add_textbox(slide, 1.0, 6.5, 11.3, 0.5,
                "每一層解決一個問題：知識 → 流程 → 分工 → 專業 → 執行 → 品管",
                font_size=18, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════
    # SLIDE 5: Layer 1 - Project Instructions
    # ══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "第 1 層：AI 的工作手冊")
    add_section_label(slide, "概念篇")
    add_layer_badge(slide, 1)
    add_page_number(slide, 5)

    # Left: description
    add_multiline(slide, 0.8, 1.6, 5.5, 2.0, [
        ("一份 Markdown 文件", WHITE, True, 24),
        ("定義 AI 在此專案中的所有規則", LIGHT_GRAY, False, 18),
    ], font_size=18)

    # Content items with icons
    items = [
        ("📖", "專業術語定義", "AI 要說你的語言", ACCENT_CYAN),
        ("📏", "工程規則與數值", "精確到小數點", ACCENT_ORANGE),
        ("🔌", "軟體 API 用法", "哪些能用、怎麼用", ACCENT_GREEN),
        ("🚫", "禁止事項", "哪些事絕不能做", ACCENT_RED),
    ]

    for i, (icon, title, desc, accent) in enumerate(items):
        y = 3.0 + i * 1.0
        add_rounded_rect(slide, 0.8, y, 5.5, 0.85, BG_CARD, accent)
        add_textbox(slide, 0.9, y + 0.05, 0.6, 0.5,
                    icon, font_size=28, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, 1.6, y + 0.05, 2.5, 0.4,
                    title, font_size=18, color=WHITE, bold=True)
        add_textbox(slide, 1.6, y + 0.42, 4.5, 0.4,
                    desc, font_size=15, color=accent)

    # Right: SOP analogy card
    add_rounded_rect(slide, 7.0, 1.6, 5.8, 5.3, BG_CARD, ACCENT_CYAN)
    add_textbox(slide, 7.2, 1.8, 5.4, 0.6,
                "類比：新進員工的 SOP 手冊", font_size=22, color=ACCENT_CYAN, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Mockup of markdown document
    add_rounded_rect(slide, 7.5, 2.6, 5.0, 3.8, BG_TERMINAL)
    md_lines = [
        ("# 專案說明書", ACCENT_ORANGE, True, 16),
        ("", WHITE, False, 6),
        ("## 術語定義", ACCENT_CYAN, True, 14),
        ("- 上構 = 1MF 以上樓層", LIGHT_GRAY, False, 13),
        ("- 下構 = B*F, 1F, BASE", LIGHT_GRAY, False, 13),
        ("", WHITE, False, 6),
        ("## 工程規則", ACCENT_CYAN, True, 14),
        ("- Rigid Zone = 0.75", LIGHT_GRAY, False, 13),
        ("- 柱保護層 = 7cm", LIGHT_GRAY, False, 13),
        ("", WHITE, False, 6),
        ("## 禁止事項", ACCENT_RED, True, 14),
        ("- 不可建立 SDL 載重", LIGHT_GRAY, False, 13),
        ("- 小梁位置不可用猜的", LIGHT_GRAY, False, 13),
    ]
    add_multiline(slide, 7.7, 2.8, 4.6, 3.5, md_lines,
                  font_size=13, font_name="Consolas", spacing=1.0)

    # Bottom analogy
    add_textbox(slide, 7.2, 6.6, 5.4, 0.4,
                "📋 還沒上工就先讀手冊 → AI 每次啟動先讀這份",
                font_size=15, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════
    # SLIDE 6: Layer 2 - Workflow Commands
    # ══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "第 2 層：一鍵啟動完整流程")
    add_section_label(slide, "概念篇")
    add_layer_badge(slide, 2)
    add_page_number(slide, 6)

    # Big button icon
    btn = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(1.8), Inches(2.5), Inches(1.5)
    )
    btn.fill.solid()
    btn.fill.fore_color.rgb = ACCENT_ORANGE
    btn.line.fill.background()
    add_textbox(slide, 1.5, 2.0, 2.5, 1.2,
                "▶ START", font_size=36, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Arrow from button to phases
    add_arrow(slide, 4.1, 2.3, 0.8, 0.5, ACCENT_ORANGE)

    # Phase flow
    phases = [
        ("Phase 1", "參數檢查", "輸入完整？", ACCENT_CYAN),
        ("Phase 2", "召集團隊", "組建 Agent", ACCENT_ORANGE),
        ("Phase 3", "依序執行", "步驟自動編排", ACCENT_GREEN),
        ("Phase 4", "品質驗收", "自動驗證結果", ACCENT_PURPLE),
    ]

    for i, (phase, title, desc, accent) in enumerate(phases):
        x = 5.1 + i * 2.1
        add_rounded_rect(slide, x, 1.8, 1.9, 2.5, BG_CARD, accent)
        add_textbox(slide, x + 0.05, 1.9, 1.8, 0.4,
                    phase, font_size=14, color=accent, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.05, 2.4, 1.8, 0.5,
                    title, font_size=18, color=WHITE, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.05, 3.0, 1.8, 0.6,
                    desc, font_size=14, color=LIGHT_GRAY,
                    alignment=PP_ALIGN.CENTER)
        if i < 3:
            add_arrow(slide, x + 1.92, 2.7, 0.15, 0.25, MID_GRAY)

    # Key points
    add_multiline(slide, 0.8, 4.8, 11.7, 2.0, [
        ("一個命令 = 一條完整的生產線", WHITE, True, 22),
        ("使用者只需一行指令，系統自動編排所有步驟", LIGHT_GRAY, False, 18),
        ("", WHITE, False, 8),
        ("類比：工廠的生產線按鈕 — 按一下，整條產線啟動", ACCENT_CYAN, False, 18),
    ], font_size=18)

    # ══════════════════════════════════════════════
    # SLIDE 7: Layer 3 - AI Agent Team
    # ══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "第 3 層：多專家分工")
    add_section_label(slide, "概念篇")
    add_layer_badge(slide, 3)
    add_page_number(slide, 7)

    # Three agent cards
    agents = [
        ("🔍", "讀圖專家", "判讀圖紙、辨識構件\n理解設計意圖", ACCENT_CYAN),
        ("📐", "定位專家", "精準定位座標\n驗證幾何合理性", ACCENT_ORANGE),
        ("📋", "整合專家", "彙整所有資訊\n產出標準化規格書", ACCENT_GREEN),
    ]

    for i, (icon, role, desc, accent) in enumerate(agents):
        x = 0.8 + i * 4.2
        add_rounded_rect(slide, x, 1.6, 3.8, 3.2, BG_CARD, accent)
        add_textbox(slide, x + 0.1, 1.8, 3.6, 0.8,
                    icon, font_size=48, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.1, 2.6, 3.6, 0.6,
                    role, font_size=26, color=accent, bold=True,
                    alignment=PP_ALIGN.CENTER)
        # Divider
        div = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x + 0.6), Inches(3.3), Inches(2.6), Inches(0.03)
        )
        div.fill.solid()
        div.fill.fore_color.rgb = accent
        div.line.fill.background()
        add_textbox(slide, x + 0.2, 3.5, 3.4, 1.2,
                    desc, font_size=17, color=LIGHT_GRAY,
                    alignment=PP_ALIGN.CENTER)

    # Communication arrows between agents
    for i in range(2):
        x = 4.7 + i * 4.2
        # Forward arrow
        add_arrow(slide, x, 2.5, 0.3, 0.25, MID_GRAY)
        # Backward arrow (using a left arrow via rotation hack - just use text)
        add_textbox(slide, x - 0.05, 3.0, 0.4, 0.3,
                    "◀▶", font_size=12, color=MID_GRAY,
                    alignment=PP_ALIGN.CENTER)

    # Key design points
    key_points = [
        ("🎭  角色定義", "每位專家只負責自己擅長的事"),
        ("🔧  工具權限", "每位專家有自己能使用的工具"),
        ("💬  直接溝通", "專家之間可互傳訊息、交叉確認"),
    ]

    for i, (title, desc) in enumerate(key_points):
        x = 0.8 + i * 4.2
        add_rounded_rect(slide, x, 5.2, 3.8, 0.85, BG_CARD)
        add_textbox(slide, x + 0.15, 5.25, 3.5, 0.4,
                    title, font_size=16, color=ACCENT_ORANGE, bold=True)
        add_textbox(slide, x + 0.15, 5.6, 3.5, 0.4,
                    desc, font_size=14, color=LIGHT_GRAY)

    # Analogy
    add_textbox(slide, 0.8, 6.4, 11.7, 0.4,
                "類比：專案小組 — 每個人有明確分工，遇到問題可以互相討論",
                font_size=16, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════
    # SLIDE 8: Layer 4 - Skill Modules
    # ══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "第 4 層：可複用的知識庫")
    add_section_label(slide, "概念篇")
    add_layer_badge(slide, 4)
    add_page_number(slide, 8)

    # Center: skill module bookshelf concept
    # Main description
    add_multiline(slide, 0.8, 1.6, 5.0, 2.0, [
        ("技能模組 = 封裝好的專業知識", WHITE, True, 22),
        ("Domain Knowledge，打包成可重複使用的模組", LIGHT_GRAY, False, 17),
    ], font_size=18)

    # Three feature cards
    features = [
        ("📚", "跨專案共享", "同一份專業知識\n可用於不同案件", ACCENT_CYAN),
        ("🔄", "隨時更新", "規則變了改一處\n全案生效", ACCENT_ORANGE),
        ("📋", "標準化", "所有 Agent 引用同一套知識\n不會各說各話", ACCENT_GREEN),
    ]

    for i, (icon, title, desc, accent) in enumerate(features):
        x = 0.8 + i * 4.2
        add_rounded_rect(slide, x, 3.2, 3.8, 2.5, BG_CARD, accent)
        add_textbox(slide, x + 0.1, 3.3, 0.7, 0.6,
                    icon, font_size=32, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.8, 3.35, 2.8, 0.5,
                    title, font_size=22, color=accent, bold=True)
        add_textbox(slide, x + 0.2, 4.1, 3.4, 1.5,
                    desc, font_size=17, color=LIGHT_GRAY,
                    alignment=PP_ALIGN.CENTER)

    # Right: visual concept - one module used by many projects
    add_rounded_rect(slide, 7.0, 1.4, 5.8, 1.5, BG_CARD, ACCENT_PURPLE)
    add_textbox(slide, 7.2, 1.5, 5.4, 0.5,
                "一份知識 → 多個專案", font_size=20, color=ACCENT_PURPLE, bold=True,
                alignment=PP_ALIGN.CENTER)

    project_names = ["專案 A", "專案 B", "專案 C", "專案 D"]
    for i, name in enumerate(project_names):
        x = 7.4 + i * 1.3
        add_textbox(slide, x, 2.2, 1.1, 0.4,
                    name, font_size=13, color=LIGHT_GRAY,
                    alignment=PP_ALIGN.CENTER)

    # Analogy
    add_textbox(slide, 0.8, 6.2, 11.7, 0.5,
                "類比：專業教科書 / 參考手冊 — 每個工程師桌上都有一本，翻開就能查",
                font_size=16, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════
    # SLIDE 9: Layer 5 - Deterministic Scripts
    # ══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "第 5 層：AI 理解，程式執行")
    add_section_label(slide, "概念篇")
    add_layer_badge(slide, 5)
    add_page_number(slide, 9)

    # Left side: AI brain
    add_rounded_rect(slide, 0.8, 1.8, 3.8, 3.5, BG_CARD, ACCENT_CYAN)
    add_textbox(slide, 0.9, 1.9, 3.6, 0.7,
                "🧠 AI 負責", font_size=28, color=ACCENT_CYAN, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_multiline(slide, 1.0, 2.7, 3.4, 2.5, [
        ("讀懂需求", WHITE, True, 20),
        ("理解自然語言", LIGHT_GRAY, False, 16),
        ("判讀圖紙", LIGHT_GRAY, False, 16),
        ("處理模糊資訊", LIGHT_GRAY, False, 16),
    ], font_size=16, alignment=PP_ALIGN.CENTER)

    # Middle: config file (bridge)
    add_arrow(slide, 4.7, 3.0, 0.6, 0.4, ACCENT_ORANGE)

    add_rounded_rect(slide, 5.4, 2.0, 2.5, 3.0, BG_CARD, ACCENT_ORANGE)
    add_textbox(slide, 5.5, 2.1, 2.3, 0.5,
                "📄 規格書", font_size=22, color=ACCENT_ORANGE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 5.5, 2.7, 2.3, 0.5,
                "設定檔", font_size=18, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    add_multiline(slide, 5.6, 3.3, 2.1, 1.5, [
        ("結構化資料", MID_GRAY, False, 14),
        ("JSON 格式", MID_GRAY, False, 14),
        ("可版控追溯", MID_GRAY, False, 14),
    ], font_size=14, alignment=PP_ALIGN.CENTER)

    add_arrow(slide, 8.0, 3.0, 0.6, 0.4, ACCENT_ORANGE)

    # Right side: Code gears
    add_rounded_rect(slide, 8.7, 1.8, 3.8, 3.5, BG_CARD, ACCENT_GREEN)
    add_textbox(slide, 8.8, 1.9, 3.6, 0.7,
                "⚙️ 程式負責", font_size=28, color=ACCENT_GREEN, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_multiline(slide, 8.9, 2.7, 3.4, 2.5, [
        ("精準動作", WHITE, True, 20),
        ("數值計算零誤差", LIGHT_GRAY, False, 16),
        ("同輸入 = 同輸出", LIGHT_GRAY, False, 16),
        ("大量重複操作", LIGHT_GRAY, False, 16),
    ], font_size=16, alignment=PP_ALIGN.CENTER)

    # Key insight
    add_rounded_rect(slide, 0.8, 5.6, 11.7, 1.3, BG_CARD, YELLOW)
    add_textbox(slide, 1.0, 5.7, 11.3, 0.5,
                "確定性 = 同樣的輸入，永遠得到同樣的結果",
                font_size=24, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1.0, 6.3, 11.3, 0.5,
                "類比：AI 是建築師，腳本是施工隊 — 建築師畫好圖，施工隊照圖施工，不會自己改設計",
                font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════
    # SLIDE 10: Layer 6 - Automated Verification
    # ══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "第 6 層：每次執行，自動品質檢查")
    add_section_label(slide, "概念篇")
    add_layer_badge(slide, 6)
    add_page_number(slide, 10)

    # Left: verification categories
    verif_items = [
        ("✅", "數值正確性", "所有參數、設定值是否正確"),
        ("✅", "完整性", "該有的東西是否都有"),
        ("✅", "一致性", "結果是否與規格書吻合"),
    ]

    add_textbox(slide, 0.8, 1.6, 6.0, 0.5,
                "每次執行完成後，自動跑品質檢測",
                font_size=22, color=WHITE, bold=True)

    for i, (icon, title, desc) in enumerate(verif_items):
        y = 2.5 + i * 1.2
        add_rounded_rect(slide, 0.8, y, 6.0, 1.0, BG_CARD, ACCENT_GREEN)
        add_textbox(slide, 1.0, y + 0.1, 0.5, 0.5,
                    icon, font_size=24, color=ACCENT_GREEN, bold=True)
        add_textbox(slide, 1.6, y + 0.1, 2.0, 0.4,
                    title, font_size=20, color=WHITE, bold=True)
        add_textbox(slide, 1.6, y + 0.5, 5.0, 0.4,
                    desc, font_size=16, color=LIGHT_GRAY)

    # Error handling
    add_rounded_rect(slide, 0.8, 5.5, 6.0, 1.0, BG_CARD, ACCENT_RED)
    add_textbox(slide, 1.0, 5.6, 0.5, 0.5,
                "❌", font_size=24, color=ACCENT_RED, bold=True)
    add_textbox(slide, 1.6, 5.6, 2.5, 0.4,
                "不通過時", font_size=20, color=ACCENT_RED, bold=True)
    add_textbox(slide, 1.6, 6.0, 5.0, 0.4,
                "精確指出哪裡有問題，不是模糊的「有錯」",
                font_size=16, color=LIGHT_GRAY)

    # Right: QC mockup
    add_rounded_rect(slide, 7.3, 1.6, 5.5, 5.3, BG_TERMINAL, ACCENT_GREEN)
    add_textbox(slide, 7.5, 1.8, 5.1, 0.5,
                "品質檢測報告", font_size=18, color=ACCENT_GREEN, bold=True,
                alignment=PP_ALIGN.CENTER)

    qc_lines = [
        ("", WHITE, False, 6),
        ("數值正確性", ACCENT_CYAN, True, 16),
        ("  ✓ 參數值在合理範圍", ACCENT_GREEN, False, 14),
        ("  ✓ 設定值與規則一致", ACCENT_GREEN, False, 14),
        ("", WHITE, False, 6),
        ("完整性", ACCENT_CYAN, True, 16),
        ("  ✓ 所有必要構件已建立", ACCENT_GREEN, False, 14),
        ("  ✓ 載重型態完整", ACCENT_GREEN, False, 14),
        ("", WHITE, False, 6),
        ("一致性", ACCENT_CYAN, True, 16),
        ("  ✓ 結果與規格書吻合", ACCENT_GREEN, False, 14),
        ("  ✓ 命名規則正確", ACCENT_GREEN, False, 14),
        ("", WHITE, False, 6),
        ("════ ALL PASSED ════", ACCENT_GREEN, True, 16),
    ]
    add_multiline(slide, 7.5, 2.3, 5.1, 4.5, qc_lines,
                  font_size=14, font_name="Consolas", spacing=1.0)

    # Analogy
    add_textbox(slide, 0.8, 7.0, 11.7, 0.4,
                "類比：出廠品管 — 每一件產品離開產線前都要過 QC，不合格不出貨",
                font_size=16, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════
    # SLIDE 11: Design Philosophy
    # ══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "各取所長 — AI 理解 + 程式執行")
    add_section_label(slide, "概念篇 · 總結")
    add_page_number(slide, 11)

    # Left: AI strengths
    add_rounded_rect(slide, 0.8, 1.6, 5.0, 3.5, BG_CARD, ACCENT_CYAN)
    add_textbox(slide, 0.9, 1.7, 4.8, 0.7,
                "🧠 AI 擅長", font_size=30, color=ACCENT_CYAN, bold=True,
                alignment=PP_ALIGN.CENTER)
    ai_strengths = [
        "理解自然語言",
        "判讀圖紙與文件",
        "處理模糊、不完整的資訊",
        "跨領域知識整合",
    ]
    for i, s in enumerate(ai_strengths):
        add_textbox(slide, 1.2, 2.5 + i * 0.6, 4.4, 0.4,
                    f"◆  {s}", font_size=18, color=WHITE)

    # Right: Code strengths
    add_rounded_rect(slide, 7.5, 1.6, 5.0, 3.5, BG_CARD, ACCENT_GREEN)
    add_textbox(slide, 7.6, 1.7, 4.8, 0.7,
                "⚙️ 程式擅長", font_size=30, color=ACCENT_GREEN, bold=True,
                alignment=PP_ALIGN.CENTER)
    code_strengths = [
        "精準執行指令",
        "數值計算零誤差",
        "同輸入 = 同輸出",
        "毫秒級大量重複操作",
    ]
    for i, s in enumerate(code_strengths):
        add_textbox(slide, 7.9, 2.5 + i * 0.6, 4.4, 0.4,
                    f"◆  {s}", font_size=18, color=WHITE)

    # Center: Bridge (config file)
    add_rounded_rect(slide, 4.5, 5.4, 4.3, 1.5, BG_CARD, ACCENT_ORANGE)
    add_textbox(slide, 4.6, 5.5, 4.1, 0.6,
                "📄 結構化設定檔（規格書）", font_size=18, color=ACCENT_ORANGE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 4.6, 6.1, 4.1, 0.6,
                "AI 產出 ←→ 程式讀取", font_size=18, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # Arrows pointing to center
    # From AI to config
    add_down_arrow(slide, 3.0, 5.2, 0.4, 0.6, ACCENT_CYAN)
    add_textbox(slide, 2.2, 5.1, 1.5, 0.3,
                "AI 產出", font_size=14, color=ACCENT_CYAN, bold=True)
    # From config to code
    add_down_arrow(slide, 9.8, 5.2, 0.4, 0.6, ACCENT_GREEN)
    add_textbox(slide, 9.6, 5.1, 1.5, 0.3,
                "程式讀取", font_size=14, color=ACCENT_GREEN, bold=True)

    # Core principle
    add_textbox(slide, 0.8, 7.0, 11.7, 0.4,
                "核心原則：讓 AI 做 AI 擅長的事，讓程式做程式擅長的事",
                font_size=20, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════
    # SLIDE 12: Application - Structural Modeling
    # ══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "案例：AI 驅動的 ETABS 自動建模")
    add_section_label(slide, "應用篇")
    add_page_number(slide, 12)

    # Pain point (top-left)
    add_rounded_rect(slide, 0.8, 1.6, 5.5, 2.0, BG_CARD, ACCENT_RED)
    add_textbox(slide, 0.9, 1.7, 5.3, 0.5,
                "痛點", font_size=22, color=ACCENT_RED, bold=True)
    add_multiline(slide, 0.9, 2.2, 5.3, 1.2, [
        ("一棟 20 層建築 → ETABS 建模需 2-3 天", WHITE, False, 17),
        ("讀圖→Grid→斷面→柱梁牆→配筋→載重→檢查，全程手動", LIGHT_GRAY, False, 15),
    ], font_size=15)

    # Solution (top-right)
    add_rounded_rect(slide, 6.8, 1.6, 5.8, 2.0, BG_CARD, ACCENT_GREEN)
    add_textbox(slide, 6.9, 1.7, 5.6, 0.5,
                "解決方案：套用六層架構", font_size=22, color=ACCENT_GREEN, bold=True)
    add_multiline(slide, 6.9, 2.2, 5.6, 1.2, [
        ("全程 15-20 分鐘，零手動", WHITE, True, 17),
        ("AI 讀圖理解 + 確定性程式碼建模 + 自動驗證", LIGHT_GRAY, False, 15),
    ], font_size=15)

    # Three-phase flow diagram
    flow_phases = [
        ("AI 讀圖", "判讀結構圖", "~3 分鐘", ACCENT_CYAN),
        ("AI 整合", "生成規格書", "~2 分鐘", ACCENT_ORANGE),
        ("程式建模", "11 步確定性執行", "~10 分鐘", ACCENT_GREEN),
    ]

    for i, (title, desc, time, accent) in enumerate(flow_phases):
        x = 0.8 + i * 4.2
        add_rounded_rect(slide, x, 4.2, 3.8, 2.5, BG_CARD, accent)
        add_textbox(slide, x + 0.1, 4.3, 3.6, 0.6,
                    title, font_size=26, color=accent, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.1, 4.9, 3.6, 0.5,
                    desc, font_size=18, color=WHITE,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.1, 5.5, 3.6, 0.6,
                    time, font_size=22, color=YELLOW, bold=True,
                    alignment=PP_ALIGN.CENTER)
        if i < 2:
            add_arrow(slide, x + 3.85, 5.0, 0.3, 0.4, MID_GRAY)

    # Big number
    add_textbox(slide, 8.5, 7.0, 4.3, 0.5,
                "2-3 天 → 15 分鐘", font_size=24, color=YELLOW, bold=True,
                alignment=PP_ALIGN.RIGHT)

    # ══════════════════════════════════════════════
    # SLIDE 13: Real Project Results
    # ══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "真實專案成果")
    add_section_label(slide, "應用篇")
    add_page_number(slide, 13)

    # Case name
    add_textbox(slide, 0.8, 1.6, 5.0, 0.6,
                "A21 捷運社宅", font_size=32, color=ACCENT_ORANGE, bold=True)

    # Data table
    case_data = [
        ("規模",     "18 層（B3F–PRF，含 3 層地下室 + 屋突）"),
        ("Grid",    "X: 5 軸  /  Y: 5 軸"),
        ("斷面種類",  "15+ 種（柱/梁/小梁/牆梁/基礎梁）"),
        ("混凝土強度", "C280–C490（分層配置）"),
        ("建模時間",  "~ 15 分鐘"),
        ("品質驗證",  "9/9 自動測試全數通過 ✓"),
    ]

    add_rounded_rect(slide, 0.8, 2.4, 6.0, 4.5, BG_CARD, ACCENT_ORANGE)
    for i, (label, val) in enumerate(case_data):
        y = 2.6 + i * 0.7
        add_textbox(slide, 1.0, y, 2.0, 0.5,
                    label, font_size=17, color=ACCENT_ORANGE, bold=True)
        add_textbox(slide, 3.0, y, 3.6, 0.5,
                    val, font_size=16, color=WHITE)

    # Right: ETABS model placeholder
    add_rounded_rect(slide, 7.3, 1.6, 5.5, 5.3, BG_CARD, ACCENT_CYAN)
    add_textbox(slide, 7.5, 3.3, 5.1, 1.0,
                "[ ETABS 3D 模型截圖 ]",
                font_size=24, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 7.5, 4.2, 5.1, 0.8,
                "請替換為實際模型渲染圖",
                font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # Demo flow at bottom
    add_textbox(slide, 0.8, 7.0, 11.7, 0.4,
                "Demo：拿到圖紙 → 一鍵啟動 → AI 讀圖 → 生成規格書 → 自動建模 → 驗證通過 → 完成",
                font_size=15, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════
    # SLIDE 14: Efficiency Comparison
    # ══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "數字會說話")
    add_section_label(slide, "應用篇")
    add_page_number(slide, 14)

    comp_items = [
        ("建模時間",   "2–3 天",       "15–20 分鐘"),
        ("人工介入",   "全程手動",      "僅審圖確認"),
        ("參數錯誤率", "偶發（難全檢）", "0（規則寫在程式碼裡）"),
        ("可重現性",   "低（因人而異）", "100%（同規格=同結果）"),
        ("品質驗證",   "人工目視",      "自動化測試"),
    ]

    # Header
    add_rounded_rect(slide, 0.8, 1.6, 11.7, 0.65, ACCENT_ORANGE)
    add_textbox(slide, 1.0, 1.65, 3.5, 0.55,
                "項目", font_size=20, color=BG_DARK, bold=True)
    add_textbox(slide, 4.8, 1.65, 3.5, 0.55,
                "🔴 傳統方式", font_size=20, color=BG_DARK, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 8.8, 1.65, 3.5, 0.55,
                "🟢 AI Agent", font_size=20, color=BG_DARK, bold=True,
                alignment=PP_ALIGN.CENTER)

    for i, (item, old, new) in enumerate(comp_items):
        y = 2.45 + i * 0.78
        bg = BG_CARD if i % 2 == 0 else BG_DARK
        add_rounded_rect(slide, 0.8, y, 11.7, 0.7, bg)
        add_textbox(slide, 1.0, y + 0.08, 3.5, 0.5,
                    item, font_size=18, color=WHITE, bold=True)
        add_textbox(slide, 4.8, y + 0.08, 3.5, 0.5,
                    old, font_size=18, color=ACCENT_RED,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, 8.8, y + 0.08, 3.5, 0.5,
                    new, font_size=18, color=ACCENT_GREEN, bold=True,
                    alignment=PP_ALIGN.CENTER)

    # Impact numbers
    add_textbox(slide, 0.8, 6.5, 5.5, 0.8,
                "省 100+ 天/年 工程師工時",
                font_size=28, color=YELLOW, bold=True)
    add_textbox(slide, 7.0, 6.5, 5.5, 0.8,
                "又快又準，每次一致",
                font_size=28, color=ACCENT_GREEN, bold=True,
                alignment=PP_ALIGN.RIGHT)

    # ══════════════════════════════════════════════
    # SLIDE 15: More Applications
    # ══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "不只是結構建模 — 架構可移植到任何場景")
    add_section_label(slide, "應用篇")
    add_page_number(slide, 15)

    # Core logic
    add_rounded_rect(slide, 0.8, 1.6, 11.7, 1.2, BG_CARD, YELLOW)
    add_textbox(slide, 1.0, 1.7, 11.3, 0.5,
                "六層架構的核心邏輯：AI 理解 + 確定性執行 + 自動驗證",
                font_size=24, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1.0, 2.2, 11.3, 0.4,
                "只要符合這個模式的工作流，都能套用",
                font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Application scenarios - 4 cards
    scenarios = [
        ("🏗️", "結構建模", "AI 讀懂結構圖", "自動建 ETABS 模型", "參數/構件品質檢查", ACCENT_CYAN),
        ("📄", "報告生成", "AI 讀懂分析結果", "自動排版出報告", "格式/內容完整性", ACCENT_ORANGE),
        ("📐", "圖紙審查", "AI 讀懂圖紙", "自動比對規範", "違規項目清單", ACCENT_GREEN),
        ("📊", "數據處理", "AI 理解數據意義", "批次清洗/轉換", "資料一致性檢查", ACCENT_PURPLE),
    ]

    for i, (icon, title, ai_part, exec_part, verify_part, accent) in enumerate(scenarios):
        x = 0.6 + i * 3.15
        add_rounded_rect(slide, x, 3.1, 2.95, 3.5, BG_CARD, accent)
        add_textbox(slide, x + 0.05, 3.2, 2.85, 0.5,
                    f"{icon}  {title}", font_size=20, color=accent, bold=True,
                    alignment=PP_ALIGN.CENTER)
        # Divider
        div = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x + 0.4), Inches(3.8), Inches(2.15), Inches(0.03)
        )
        div.fill.solid()
        div.fill.fore_color.rgb = accent
        div.line.fill.background()

        add_multiline(slide, x + 0.15, 4.0, 2.65, 2.5, [
            ("AI 理解", accent, True, 14),
            (ai_part, LIGHT_GRAY, False, 13),
            ("", WHITE, False, 6),
            ("確定性執行", accent, True, 14),
            (exec_part, LIGHT_GRAY, False, 13),
            ("", WHITE, False, 6),
            ("自動驗證", accent, True, 14),
            (verify_part, LIGHT_GRAY, False, 13),
        ], font_size=13)

    # Key message
    add_textbox(slide, 0.8, 6.8, 11.7, 0.5,
                "關鍵：架構是通用的，專業知識是可抽換的",
                font_size=20, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════
    # SLIDE 16: Closing + CTA
    # ══════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Top bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_ORANGE
    bar.line.fill.background()

    # Title
    add_textbox(slide, 1.0, 1.5, 11.3, 1.0,
                "讓 AI 幫你做好專業工作",
                font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Divider
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.0), Inches(2.7), Inches(3.333), Inches(0.04)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = ACCENT_ORANGE
    divider.line.fill.background()

    # CTA points
    cta_items = [
        ("✅  Pilot 合作", "提供一份實際案例，我們現場展示成果"),
        ("✅  客製化服務", "根據您的工作流程、規範、標準量身打造"),
        ("✅  持續進化", "每個合作案例都讓系統更強大"),
    ]

    for i, (title, desc) in enumerate(cta_items):
        y = 3.2 + i * 0.9
        add_textbox(slide, 2.0, y, 5.0, 0.4,
                    title, font_size=24, color=ACCENT_GREEN, bold=True)
        add_textbox(slide, 7.0, y + 0.05, 5.0, 0.4,
                    desc, font_size=18, color=LIGHT_GRAY)

    # Tagline
    add_textbox(slide, 1.0, 5.3, 11.3, 0.5,
                "一起定義 AI 時代的專業工作方式",
                font_size=22, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

    # Contact
    add_rounded_rect(slide, 3.5, 5.9, 6.3, 1.0, BG_CARD, ACCENT_ORANGE)
    add_textbox(slide, 3.7, 6.0, 5.9, 0.8,
                "聯絡我們：[ Email / 電話 / Line ]\n[ 公司名稱 · Logo ]",
                font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Bottom bar
    bar2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.42), Inches(13.333), Inches(0.08)
    )
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = ACCENT_ORANGE
    bar2.line.fill.background()

    add_page_number(slide, 16)

    # ── Save ──
    output_path = os.path.join(os.path.dirname(__file__),
                               "AI_ETABS_自動建模系統.pptx")
    prs.save(output_path)
    print(f"PPT saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    create_ppt()
