"""Tests for pptx_to_elements helpers (vertex counting)."""

import sys
import os
import pytest
from lxml import etree

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from golden_scripts.tools.pptx_to_elements import (
    _get_freeform_vertex_count,
    _AbsShape,
    _iter_slide_shapes,
    COLUMN_REQUIRED_VERTICES,
    FLOOR_LABEL_RE,
    expand_floor_range,
    format_scan_floors_output,
    parse_page_floors,
)

# ─── XML helpers ─────────────────────────────────────────────────────────────

_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _make_shape_with_vertices(move_count, line_count):
    """Build a minimal mock shape whose _element has custGeom with the given vertex counts."""
    moves = "".join(f'<a:moveTo xmlns:a="{_NS}"><a:pt x="0" y="0"/></a:moveTo>' for _ in range(move_count))
    lines = "".join(f'<a:lnTo xmlns:a="{_NS}"><a:pt x="100" y="100"/></a:lnTo>' for _ in range(line_count))

    xml = f"""
    <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="{_NS}">
      <p:spPr>
        <a:custGeom>
          <a:pathLst>
            <a:path w="100" h="100">
              {moves}
              {lines}
              <a:close/>
            </a:path>
          </a:pathLst>
        </a:custGeom>
      </p:spPr>
    </p:sp>
    """
    el = etree.fromstring(xml.encode())

    class MockShape:
        _element = el

    return MockShape()


def _make_shape_no_custgeom():
    """Build a mock shape with no custGeom (e.g. a preset shape)."""
    xml = f"""
    <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="{_NS}">
      <p:spPr>
        <a:prstGeom prst="rect"/>
      </p:spPr>
    </p:sp>
    """
    el = etree.fromstring(xml.encode())

    class MockShape:
        _element = el

    return MockShape()


# ─── Tests ───────────────────────────────────────────────────────────────────


class TestGetFreeformVertexCount:
    """Tests for _get_freeform_vertex_count."""

    def test_rectangle_4_vertices(self):
        """A 4-vertex freeform (1 moveTo + 3 lnTo) should return 4."""
        shape = _make_shape_with_vertices(move_count=1, line_count=3)
        assert _get_freeform_vertex_count(shape) == 4

    def test_triangle_3_vertices(self):
        """A 3-vertex freeform (1 moveTo + 2 lnTo) should return 3 (rejected as column)."""
        shape = _make_shape_with_vertices(move_count=1, line_count=2)
        result = _get_freeform_vertex_count(shape)
        assert result == 3
        assert result != COLUMN_REQUIRED_VERTICES

    def test_irregular_polygon_6_vertices(self):
        """A 6-vertex freeform should return 6 (rejected as column)."""
        shape = _make_shape_with_vertices(move_count=1, line_count=5)
        result = _get_freeform_vertex_count(shape)
        assert result == 6
        assert result != COLUMN_REQUIRED_VERTICES

    def test_no_custgeom_returns_none(self):
        """Shape without custGeom should return None (accepted as fallback)."""
        shape = _make_shape_no_custgeom()
        assert _get_freeform_vertex_count(shape) is None

    def test_no_element_returns_none(self):
        """Shape without _element attribute should return None."""

        class NoElement:
            pass

        assert _get_freeform_vertex_count(NoElement()) is None

    def test_constant_value(self):
        """COLUMN_REQUIRED_VERTICES should be 4."""
        assert COLUMN_REQUIRED_VERTICES == 4


# ─── Floor label regex tests ─────────────────────────────────────────────────


class TestFloorLabelRegex:
    """Tests for FLOOR_LABEL_RE — Unicode word boundary fix."""

    def test_chinese_suffix_no_space(self):
        """B1F~B2F結構配置圖 → matches B1F~B2F (Chinese after label, no space)."""
        assert FLOOR_LABEL_RE.findall("B1F~B2F結構配置圖") == ["B1F~B2F"]

    def test_chinese_comma_separator(self):
        """B1F、B2F結構配置圖 → matches B1F and B2F individually."""
        assert FLOOR_LABEL_RE.findall("B1F、B2F結構配置圖") == ["B1F", "B2F"]

    def test_fullwidth_tilde_range(self):
        """B1F～B2F結構配置圖 → matches B1F～B2F (full-width tilde)."""
        assert FLOOR_LABEL_RE.findall("B1F\uff5eB2F結構配置圖") == ["B1F\uff5eB2F"]

    def test_negative_C490F(self):
        """C490F → no match (regression guard for section names like C490F)."""
        assert FLOOR_LABEL_RE.findall("C490F") == []

    def test_space_still_works(self):
        """B1F~B2F 結構配置圖 → matches B1F~B2F (space before Chinese)."""
        assert FLOOR_LABEL_RE.findall("B1F~B2F 結構配置圖") == ["B1F~B2F"]


class TestExpandFloorRange:
    """Tests for expand_floor_range — separator normalization."""

    def test_fullwidth_tilde(self):
        """B1F～B2F → ['B1F', 'B2F'] (full-width tilde normalized)."""
        assert expand_floor_range("B1F\uff5eB2F") == ["B1F", "B2F"]

    def test_dash_separator(self):
        """B1F-B2F → ['B1F', 'B2F'] (dash normalized to tilde)."""
        assert expand_floor_range("B1F-B2F") == ["B1F", "B2F"]

    def test_normal_range(self):
        """3F~14F → ['3F', '4F', ..., '14F'] (existing behavior preserved)."""
        result = expand_floor_range("3F~14F")
        assert result == [f"{i}F" for i in range(3, 15)]


class TestFormatParseRoundtrip:
    """Tests for format_scan_floors_output → parse_page_floors round-trip."""

    def test_multi_label_slide_roundtrip(self):
        """Two labels on one slide survive format → parse round-trip."""
        detected = {
            3: [
                {"label": "B1F", "confidence": "high", "source_text": "B1F、B2F結構配置圖", "position": "title"},
                {"label": "B2F", "confidence": "high", "source_text": "B1F、B2F結構配置圖", "position": "title"},
            ]
        }
        output = format_scan_floors_output(detected)
        # Extract the suggested --page-floors line
        for line in output.split("\n"):
            line = line.strip()
            if line.startswith('"') and "=" in line:
                page_floors_str = line.strip('"')
                break
        parsed = parse_page_floors(page_floors_str)
        assert 3 in parsed
        assert "B1F" in parsed[3]
        assert "B2F" in parsed[3]


# ─── AbsShape group dimension scaling tests ──────────────────────────────────


from unittest.mock import MagicMock, PropertyMock
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _make_mock_shape(left, top, width, height, shape_type=MSO_SHAPE_TYPE.FREEFORM):
    """Create a mock shape with position/dimension attributes."""
    s = MagicMock()
    s.left = left
    s.top = top
    s.width = width
    s.height = height
    s.shape_type = shape_type
    return s


class TestAbsShapeGroupScaling:
    """Tests for _AbsShape width/height scaling in group transforms."""

    def test_absshape_width_height_override(self):
        """_AbsShape stores explicit width/height when provided."""
        raw = _make_mock_shape(100, 200, 300, 400)
        proxy = _AbsShape(raw, 500, 600, abs_width=900, abs_height=1200)
        assert proxy.width == 900
        assert proxy.height == 1200
        assert proxy.left == 500
        assert proxy.top == 600

    def test_absshape_width_height_default(self):
        """_AbsShape falls back to raw shape's width/height when not provided."""
        raw = _make_mock_shape(100, 200, 300, 400)
        proxy = _AbsShape(raw, 500, 600)
        assert proxy.width == 300
        assert proxy.height == 400

    def test_iter_slide_shapes_scales_grouped_dimensions(self):
        """Grouped shapes get scaled width/height from group transform."""
        # Child shape: 100x200 at (50, 50) in child coord space
        child = _make_mock_shape(50, 50, 100, 200, MSO_SHAPE_TYPE.FREEFORM)

        # Group transform: child space 1000x1000, output space 2000x4000
        # → scale_x = 2.0, scale_y = 4.0
        group = MagicMock()
        group.shape_type = MSO_SHAPE_TYPE.GROUP
        group.shapes = [child]

        # Build group xfrm XML: off=(0,0), ext=(2000,4000), chOff=(0,0), chExt=(1000,1000)
        _NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
        grp_xml = f"""
        <p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                  xmlns:a="{_NS_A}">
          <p:grpSpPr>
            <a:xfrm>
              <a:off x="0" y="0"/>
              <a:ext cx="2000" cy="4000"/>
              <a:chOff x="0" y="0"/>
              <a:chExt cx="1000" cy="1000"/>
            </a:xfrm>
          </p:grpSpPr>
        </p:grpSp>
        """
        group._element = etree.fromstring(grp_xml.encode())
        group.left = 0
        group.top = 0
        group.width = 2000
        group.height = 4000

        # Mock container that holds the group
        container = MagicMock()
        container.shapes = [group]

        results = list(_iter_slide_shapes(container))
        assert len(results) == 1
        proxy = results[0]
        # Position: (0 + (50-0)*2000/1000, 0 + (50-0)*4000/1000) = (100, 200)
        assert proxy.left == 100
        assert proxy.top == 200
        # Dimensions: (100*2000/1000, 200*4000/1000) = (200, 800)
        assert proxy.width == 200
        assert proxy.height == 800

    def test_iter_slide_shapes_identity_group_no_change(self):
        """Identity group transform preserves original dimensions."""
        child = _make_mock_shape(50, 50, 300, 400, MSO_SHAPE_TYPE.FREEFORM)

        group = MagicMock()
        group.shape_type = MSO_SHAPE_TYPE.GROUP
        group.shapes = [child]

        # Identity: ext == chExt
        _NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
        grp_xml = f"""
        <p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                  xmlns:a="{_NS_A}">
          <p:grpSpPr>
            <a:xfrm>
              <a:off x="100" y="200"/>
              <a:ext cx="5000" cy="5000"/>
              <a:chOff x="0" y="0"/>
              <a:chExt cx="5000" cy="5000"/>
            </a:xfrm>
          </p:grpSpPr>
        </p:grpSp>
        """
        group._element = etree.fromstring(grp_xml.encode())
        group.left = 100
        group.top = 200
        group.width = 5000
        group.height = 5000

        container = MagicMock()
        container.shapes = [group]

        results = list(_iter_slide_shapes(container))
        assert len(results) == 1
        proxy = results[0]
        # Identity scale → dimensions unchanged
        assert proxy.width == 300
        assert proxy.height == 400
